from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# ── Color Palette ──────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x0D, 0x1A)   # very dark navy
C_PANEL    = RGBColor(0x1A, 0x1A, 0x2E)   # dark panel
C_ACCENT   = RGBColor(0xFF, 0x4D, 0x00)   # Palo Alto orange
C_ACCENT2  = RGBColor(0x00, 0xC8, 0xFF)   # cyan
C_GREEN    = RGBColor(0x00, 0xE5, 0x76)   # success green
C_RED      = RGBColor(0xFF, 0x3B, 0x3B)   # fail red
C_YELLOW   = RGBColor(0xFF, 0xC4, 0x00)   # warning yellow
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY     = RGBColor(0xAA, 0xAA, 0xBB)
C_ROW_ODD  = RGBColor(0x16, 0x16, 0x28)
C_ROW_EVEN = RGBColor(0x1E, 0x1E, 0x38)
C_HDR      = RGBColor(0xFF, 0x4D, 0x00)

blank_layout = prs.slide_layouts[6]  # blank


def set_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=Pt(14), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 16, 9, fill=C_BG)
    add_rect(slide, 0, 0, 16, 1.3, fill=C_PANEL)
    add_rect(slide, 0, 1.28, 16, 0.06, fill=C_ACCENT)
    add_text(slide, title, 0.5, 0.18, 12, 0.7, size=Pt(32), bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.82, 14, 0.4, size=Pt(13), color=C_GRAY)
    add_text(slide, "CONFIDENTIAL  |  Red Team Exercise Report  |  2026", 0.5, 8.6, 15, 0.35,
             size=Pt(9), color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# Slide 1 – Cover
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_rect(slide, 0, 0, 16, 9, fill=C_BG)
add_rect(slide, 0, 3.5, 16, 0.08, fill=C_ACCENT)
add_rect(slide, 0, 5.2, 16, 0.04, fill=C_ACCENT2)

add_text(slide, "RED TEAM ATTACK EXERCISE", 1, 1.2, 14, 1.0,
         size=Pt(42), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "전체 공격 시나리오 결과 보고서", 1, 2.4, 14, 0.7,
         size=Pt(24), bold=False, color=C_ACCENT2, align=PP_ALIGN.CENTER)

add_text(slide, "Ubuntu 14.04  ·  172.16.157.190  (XDR 미설치)", 1, 3.8, 14, 0.5,
         size=Pt(15), color=C_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Windows 11  ·  172.16.157.187  (Cortex XDR ACTIVE)", 1, 4.3, 14, 0.5,
         size=Pt(15), color=C_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "Attacker  ·  Kali Linux  ·  172.16.157.177", 1, 5.0, 14, 0.45,
         size=Pt(13), color=RGBColor(0x88, 0x88, 0x99), align=PP_ALIGN.CENTER)

add_text(slide, "2026-05-29", 1, 6.5, 14, 0.45,
         size=Pt(13), color=C_ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "CONFIDENTIAL", 1, 7.1, 14, 0.45,
         size=Pt(11), color=C_RED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# Slide 2 – Test Environment
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "테스트 환경 구성", "Lab Topology & Target Overview")

boxes = [
    ("ATTACKER", "Kali Linux\n172.16.157.177", C_RED),
    ("TARGET A", "Ubuntu 14.04\n172.16.157.190\nXDR 미설치", C_YELLOW),
    ("TARGET B", "Windows 11\n172.16.157.187\nCortex XDR ACTIVE", C_ACCENT),
]
for i, (label, body, color) in enumerate(boxes):
    x = 1.2 + i * 4.8
    add_rect(slide, x, 1.8, 4.2, 3.2, fill=C_PANEL, line=color, line_w=Pt(2))
    add_rect(slide, x, 1.8, 4.2, 0.55, fill=color)
    add_text(slide, label, x + 0.1, 1.85, 4.0, 0.45, size=Pt(16), bold=True,
             color=C_BG if color == C_YELLOW else C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.15, 2.45, 3.9, 2.4, size=Pt(14), color=C_WHITE,
             align=PP_ALIGN.CENTER)

# Services list
services = [
    "Ubuntu Services: SSH(22) · FTP(21) · HTTP/DVWA(80) · WebGoat(8080) · Mutillidae",
    "Windows Services: WinRM(5985) · SMB(445) · DVWA · RDP(3389)",
    "Kali Tools: nmap · nikto · sqlmap · hydra · metasploit · evil-winrm · john",
]
for i, s in enumerate(services):
    add_text(slide, s, 0.5, 5.4 + i * 0.55, 15, 0.45, size=Pt(12), color=C_GRAY)


# ══════════════════════════════════════════════════════════════
# Slide 3 – Attack Summary Stats
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "공격 요약 통계", "Attack Summary Statistics")

stats = [
    ("총 공격 기법", "44", C_ACCENT),
    ("성공", "40", C_GREEN),
    ("실패", "4", C_RED),
    ("MITRE 전술", "11", C_ACCENT2),
    ("XDR 감지 예상", "15+", C_YELLOW),
]
for i, (label, val, color) in enumerate(stats):
    x = 0.5 + i * 3.0
    add_rect(slide, x, 2.0, 2.7, 2.8, fill=C_PANEL, line=color, line_w=Pt(2))
    add_text(slide, val, x + 0.1, 2.35, 2.5, 1.3, size=Pt(52), bold=True,
             color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, 3.75, 2.5, 0.7, size=Pt(13),
             color=C_WHITE, align=PP_ALIGN.CENTER)

# MITRE Tactics Bar
tactics = ["Recon","Initial Access","Execution","Persistence","Priv Esc","Defense Evasion",
           "Credential Access","Discovery","Collection","Lateral Move","Impact"]
add_text(slide, "MITRE ATT&CK 전술 커버리지", 0.5, 5.2, 15, 0.4, size=Pt(13),
         bold=True, color=C_WHITE)
for i, t in enumerate(tactics):
    x = 0.5 + i * 1.37
    add_rect(slide, x, 5.7, 1.22, 0.5, fill=C_ACCENT, line=None)
    add_text(slide, t, x + 0.02, 5.72, 1.18, 0.46, size=Pt(7.5), bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# Helper – draw table
# ══════════════════════════════════════════════════════════════
def draw_table(slide, rows, col_widths, top, row_h=0.38, header=True, left=0.3):
    y = top
    for ri, row in enumerate(rows):
        x = left
        is_header = (ri == 0 and header)
        bg = C_HDR if is_header else (C_ROW_ODD if ri % 2 == 1 else C_ROW_EVEN)
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, y, cw - 0.03, row_h, fill=bg)
            color = C_BG if is_header else (
                C_GREEN if cell == "성공" else
                C_RED   if cell == "실패" else
                C_YELLOW if "부분" in str(cell) else
                C_WHITE
            )
            fsize = Pt(9.5) if is_header else Pt(8.5)
            bold = is_header
            add_text(slide, str(cell), x + 0.05, y + 0.03, cw - 0.12, row_h - 0.06,
                     size=fsize, bold=bold, color=color, align=PP_ALIGN.LEFT)
            x += cw
        y += row_h
    return y


# ══════════════════════════════════════════════════════════════
# Slide 4 – Ubuntu Attack Table (Part 1)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Ubuntu 172.16.157.190 — 공격 순서 (1/2)", "XDR 미설치 | 전체 공격 성공")

ubuntu_rows_1 = [
    ["#", "공격 기법", "MITRE", "도구", "결과"],
    ["1", "nmap 포트 스캔", "T1046", "nmap", "성공"],
    ["2", "nikto 웹 스캔", "T1595", "nikto", "성공"],
    ["3", "Hydra SSH 브루트포스", "T1110.001", "hydra", "성공"],
    ["4", "FTP 익명 로그인", "T1078", "ftp/curl", "성공"],
    ["5", "vsftpd 2.3.4 백도어 (port 6200)", "T1190", "nc", "실패"],
    ["6", "DVWA PHP 웹셸 업로드", "T1505.003", "curl/browser", "성공"],
    ["7", "리버스 셸 실행", "T1059.004", "msfvenom", "성공"],
    ["8", "CVE-2021-4034 Pwnkit", "T1068", "gcc/exploit", "실패"],
    ["9", "webadmin sudo → root", "T1548.003", "sudo/ssh", "성공"],
    ["10", "/etc/shadow 덤프", "T1003.008", "cat", "성공"],
    ["11", "John the Ripper 크랙", "T1110.002", "john", "성공"],
]
col_w = [0.4, 3.8, 1.4, 2.2, 1.0]
draw_table(slide, ubuntu_rows_1, col_w, top=1.5, row_h=0.38)

# ══════════════════════════════════════════════════════════════
# Slide 5 – Ubuntu Attack Table (Part 2)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Ubuntu 172.16.157.190 — 공격 순서 (2/2)", "Web Attacks · Persistence · Credential Access")

ubuntu_rows_2 = [
    ["#", "공격 기법", "MITRE", "도구", "결과"],
    ["12", "MySQL root 무인증 접근", "T1078.001", "mysql cli", "성공"],
    ["13", "SQLmap DB 전체 덤프", "T1190", "sqlmap", "성공"],
    ["14", "SQLi (Union/Error/Blind/Time)", "T1190", "sqlmap/manual", "성공"],
    ["15", "XSS (Reflected/Stored)", "T1059.007", "browser/curl", "성공"],
    ["16", "LFI (/etc/passwd)", "T1083", "curl", "성공"],
    ["17", "IDOR (타 사용자 프로필)", "T1087", "browser", "성공"],
    ["18", "약한 세션 ID (Mutillidae)", "T1539", "browser", "성공"],
    ["19", "WebGoat SQLi + H2 DB 추출", "T1190", "curl/ftp", "성공"],
    ["20", "cron 백도어 등록", "T1053.003", "crontab", "성공"],
    ["21", "SUID bash 백도어", "T1548.001", "chmod", "성공"],
    ["22", "SSH authorized_keys 추가", "T1098.004", "ssh-keygen", "성공"],
    ["23", "bash_history/auth.log 삭제", "T1070.003", "rm/truncate", "성공"],
]
col_w = [0.4, 3.8, 1.4, 2.2, 1.0]
draw_table(slide, ubuntu_rows_2, col_w, top=1.5, row_h=0.38)


# ══════════════════════════════════════════════════════════════
# Slide 6 – Windows Attack Table (Part 1)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Windows 172.16.157.187 — 공격 순서 (1/2)", "Cortex XDR ACTIVE | XDR Self-Protection 확인됨")

win_rows_1 = [
    ["#", "공격 기법", "MITRE", "도구", "XDR 반응", "결과"],
    ["1", "nmap SMB/WinRM 스캔", "T1046", "nmap", "-", "성공"],
    ["2", "Evil-WinRM 원격 접속", "T1021.006", "evil-winrm", "감지 가능", "성공"],
    ["3", "LSASS 덤프 (comsvcs.dll)", "T1003.001", "rundll32", "감지", "성공"],
    ["4", "SAM 덤프 (reg save)", "T1003.002", "reg.exe", "감지", "부분 성공"],
    ["5", "Meterpreter 실행", "T1059.001", "msfvenom", "감지", "성공"],
    ["6", "Bind 셸 (port 4444)", "T1059", "msfvenom", "프로세스 킬", "성공→차단"],
    ["7", "AMSI 우회 (리플렉션)", "T1562.001", "PowerShell", "우회 성공", "성공"],
    ["8", "Base64 인코딩 PowerShell", "T1027", "PowerShell", "감지 가능", "성공"],
    ["9", "In-memory shellcode 시뮬", "T1620", "PowerShell", "감지 가능", "성공"],
    ["10", "UAC 우회 (fodhelper)", "T1548.002", "fodhelper.exe", "감지", "성공"],
    ["11", "XDR 서비스 종료 시도", "T1562", "sc.exe", "액세스 거부", "실패"],
]
col_w = [0.4, 3.4, 1.4, 2.1, 1.7, 1.3]
draw_table(slide, win_rows_1, col_w, top=1.5, row_h=0.38)


# ══════════════════════════════════════════════════════════════
# Slide 7 – Windows Attack Table (Part 2)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Windows 172.16.157.187 — 공격 순서 (2/2)", "Persistence · Defense Evasion · Discovery")

win_rows_2 = [
    ["#", "공격 기법", "MITRE", "도구", "XDR 반응", "결과"],
    ["12", "스케줄 작업 등록 (schtasks)", "T1053.005", "schtasks.exe", "감지 가능", "성공"],
    ["13", "레지스트리 Run 키 등록", "T1547.001", "reg.exe", "감지 가능", "성공"],
    ["14", "악성 서비스 등록 (sc.exe)", "T1543.003", "sc.exe", "감지 가능", "성공"],
    ["15", "백도어 계정 + Administrators", "T1136.001", "net user", "감지 가능", "성공"],
    ["16", "Defender 제외 경로 추가", "T1562.001", "PowerShell", "감지 가능", "성공"],
    ["17", "WDigest 캐싱 활성화", "T1112", "reg.exe", "감지 가능", "성공"],
    ["18", "이벤트 로그 삭제 (Security+System)", "T1070.001", "wevtutil", "감지", "성공"],
    ["19", "방화벽 비활성화", "T1562.004", "netsh", "감지 가능", "성공"],
    ["20", "certutil LOLBin 활용", "T1218", "certutil.exe", "감지 가능", "성공"],
    ["21", "bitsadmin LOLBin 활용", "T1218", "bitsadmin.exe", "감지 가능", "성공"],
    ["22", "내부망 네트워크 스캔", "T1018", "PowerShell/nmap", "-", "성공"],
]
col_w = [0.4, 3.4, 1.4, 2.1, 1.7, 1.3]
draw_table(slide, win_rows_2, col_w, top=1.5, row_h=0.38)


# ══════════════════════════════════════════════════════════════
# Slide 8 – XDR Detection Highlight
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Cortex XDR 탐지 하이라이트", "XDR Self-Protection & Expected Alerts")

# Left panel – detected
add_rect(slide, 0.3, 1.5, 7.4, 5.8, fill=C_PANEL, line=C_GREEN, line_w=Pt(1.5))
add_rect(slide, 0.3, 1.5, 7.4, 0.55, fill=C_GREEN)
add_text(slide, "XDR 탐지 예상 이벤트", 0.4, 1.55, 7.2, 0.45,
         size=Pt(14), bold=True, color=C_BG, align=PP_ALIGN.CENTER)

detected = [
    "● LSASS 메모리 접근 (comsvcs.dll MiniDump)",
    "● 의심 PE 실행 (svchost_update.exe / bind.exe)",
    "● bind.exe 프로세스 강제 종료 (XDR 차단)",
    "● AMSI 우회 시도 (PowerShell 리플렉션)",
    "● 이벤트 로그 삭제 (wevtutil cl Security)",
    "● schtasks / sc.exe 악성 등록",
    "● 백도어 계정 생성 (net user)",
    "● certutil / bitsadmin 네트워크 요청",
    "● 레지스트리 WDigest 수정",
    "● Defender 제외 경로 추가",
    "● UAC 우회 (fodhelper.exe)",
    "● Meterpreter 역방향 연결 시도",
]
for i, item in enumerate(detected):
    add_text(slide, item, 0.5, 2.15 + i * 0.38, 7.0, 0.36,
             size=Pt(10.5), color=C_GREEN if "차단" in item else C_WHITE)

# Right panel – XDR self-protection
add_rect(slide, 8.1, 1.5, 7.6, 5.8, fill=C_PANEL, line=C_ACCENT, line_w=Pt(1.5))
add_rect(slide, 8.1, 1.5, 7.6, 0.55, fill=C_ACCENT)
add_text(slide, "XDR Self-Protection 확인", 8.2, 1.55, 7.4, 0.45,
         size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

sp_items = [
    ("sc.exe stop cyserver", "액세스가 거부되었습니다"),
    ("Stop-Service cyserver", "Access Denied"),
    ("taskkill /PID <xdr>", "액세스가 거부되었습니다"),
    ("net stop cyserver", "거부됨"),
]
add_text(slide, "Administrator 권한으로 시도한 XDR 종료 — 전부 실패:", 8.2, 2.2, 7.3, 0.4,
         size=Pt(11), color=C_GRAY)
for i, (cmd, result) in enumerate(sp_items):
    y = 2.75 + i * 0.9
    add_rect(slide, 8.3, y, 7.2, 0.35, fill=RGBColor(0x0D, 0x0D, 0x0D))
    add_text(slide, f"> {cmd}", 8.4, y + 0.02, 7.0, 0.32, size=Pt(10), color=C_YELLOW)
    add_text(slide, f"  {result}", 8.4, y + 0.34, 7.0, 0.35, size=Pt(10), color=C_RED)

add_text(slide, "★  XDR Self-Protection 정상 동작 확인", 8.2, 6.5, 7.2, 0.5,
         size=Pt(13), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# Slide 9 – XSIAM Verification Query
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "XSIAM 콘솔 검증 XQL", "XSIAM Console Verification Queries")

add_rect(slide, 0.5, 1.6, 15, 0.45, fill=C_PANEL)
add_text(slide, "XSIAM 콘솔 → Investigate → XQL Query 에서 아래 쿼리 실행", 0.6, 1.65, 14, 0.38,
         size=Pt(12), color=C_GRAY)

queries = [
    ("전체 알림 조회",
     'dataset = alerts\n| filter host_name = "Song-Test-oc"\n| fields alert_name, description, action_pretty, actor_process_image_name, event_timestamp\n| sort desc event_timestamp'),
    ("LSASS 접근 감지",
     'dataset = xdr_data\n| filter event_type = "PROCESS"\n   and actor_process_image_name ~= "lsass"\n| fields agent_hostname, actor_process_image_name, action_process_image_name'),
    ("이벤트 로그 삭제 감지",
     'dataset = xdr_data\n| filter event_type = "PROCESS"\n   and action_process_image_name ~= "wevtutil"\n| fields agent_hostname, action_process_image_name, action_process_image_command_line'),
]

y = 2.2
for title, q in queries:
    add_text(slide, title, 0.6, y, 14.5, 0.38, size=Pt(12), bold=True, color=C_ACCENT2)
    add_rect(slide, 0.5, y + 0.38, 15, len(q.split('\n')) * 0.32 + 0.15, fill=RGBColor(0x0A, 0x0A, 0x18))
    add_text(slide, q, 0.7, y + 0.42, 14.5, len(q.split('\n')) * 0.32,
             size=Pt(10), color=C_GREEN)
    y += len(q.split('\n')) * 0.32 + 0.75


# ══════════════════════════════════════════════════════════════
# Slide 10 – Key Findings & Recommendations
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "주요 발견 사항 & 권고사항", "Key Findings & Recommendations")

findings = [
    ("CRITICAL", "Ubuntu — XDR 미설치로 공격 전 과정 탐지 불가", C_RED),
    ("HIGH", "LSASS 메모리 덤프 성공 → 자격증명 탈취 위험", C_ACCENT),
    ("HIGH", "AMSI 우회 성공 → 메모리 기반 공격 탐지 공백", C_ACCENT),
    ("HIGH", "이벤트 로그 전체 삭제 성공 → 포렌식 증거 인멸", C_ACCENT),
    ("MEDIUM", "웹 취약점 다수 (SQLi/XSS/LFI/IDOR) — DVWA/WebGoat", C_YELLOW),
    ("MEDIUM", "백도어 계정 생성 + Administrators 그룹 추가 성공", C_YELLOW),
    ("INFO", "XDR Self-Protection 정상 동작 확인 (서비스 종료 불가)", C_GREEN),
    ("INFO", "Bind 셸 연결 후 XDR가 프로세스 자동 종료 확인", C_GREEN),
]

for i, (level, finding, color) in enumerate(findings):
    y = 1.6 + i * 0.78
    add_rect(slide, 0.4, y, 1.4, 0.55, fill=color)
    add_text(slide, level, 0.42, y + 0.05, 1.36, 0.45, size=Pt(11), bold=True,
             color=C_BG if color == C_YELLOW else C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.85, y, 13.7, 0.55, fill=C_PANEL)
    add_text(slide, finding, 1.95, y + 0.07, 13.5, 0.42, size=Pt(12), color=C_WHITE)

add_text(slide, "권고: Ubuntu 호스트에 Cortex XDR 에이전트 즉시 배포 필요", 0.4, 8.15, 15.2, 0.45,
         size=Pt(11), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════
out = "/Users/songhyeonsu/Desktop/python/RedTeam_Attack_Report.pptx"
prs.save(out)
print(f"Saved: {out}")
